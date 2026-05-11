"""
Export Pipeline — مذكرتي Pro v17

This is the single entry point for PPTX generation.
Deterministic. No threads. No subprocess. No temp files managed outside.
Validates output before returning bytes.
"""
from __future__ import annotations

import io
import logging
import shutil
import subprocess
import sys
import os
from dataclasses import dataclass

from pptx import Presentation
from pptx.util import Cm

from core.models import PresentationRequest
from core.themes import get_theme, Theme
from engine.slides import (
    set_font,
    make_cover, make_intro, make_plan, make_problem,
    make_objectives, make_importance, make_methodology,
    make_stats, make_results, make_conclusion,
    make_recommendations, make_future, make_references,
    make_final,
)

log = logging.getLogger(__name__)

# Slide dimensions
W_CM, H_CM = 33.867, 19.05

# Minimum valid PPTX size (bytes)
MIN_VALID_SIZE = 5_000


# ── Font detection ────────────────────────────────────────────────────
def _detect_arabic_font() -> str:
    """Return best available Arabic font name."""
    candidates = ["Cairo", "Amiri", "Tahoma", "Arial Unicode MS", "Calibri"]
    
    # Try fc-list
    if shutil.which("fc-list"):
        try:
            out = subprocess.run(
                ["fc-list", "--format=%{family}\n"],
                capture_output=True, text=True, timeout=5
            ).stdout.lower()
            for font in candidates:
                if font.lower() in out:
                    log.info(f"Arabic font detected: {font}")
                    return font
        except Exception:
            pass
    
    # Try common font directories
    font_dirs = [
        "/usr/share/fonts", "/usr/local/share/fonts",
        os.path.expanduser("~/.fonts"), "/tmp/fonts",
        "C:/Windows/Fonts", "/Library/Fonts",
    ]
    for font in candidates[:3]:  # Only check Arabic-capable fonts
        for d in font_dirs:
            if not os.path.isdir(d):
                continue
            for root, _, files in os.walk(d):
                for f in files:
                    if font.lower() in f.lower() and f.lower().endswith((".ttf", ".otf")):
                        log.info(f"Arabic font found on disk: {font}")
                        return font
    
    log.warning("No Arabic font found; using Calibri fallback")
    return "Calibri"


@dataclass
class ExportResult:
    success: bool
    data: bytes = b""
    slide_count: int = 0
    font_used: str = ""
    error: str = ""


class PPTXExportPipeline:
    """
    Single-responsibility export pipeline.
    
    Usage:
        pipeline = PPTXExportPipeline()
        result = pipeline.build(request)
        if result.success:
            pptx_bytes = result.data
    """

    def __init__(self):
        self._font = _detect_arabic_font()
        set_font(self._font)
        log.info(f"Export pipeline initialized | font={self._font}")

    def build(self, req: PresentationRequest) -> ExportResult:
        """
        Build a PPTX from a validated request.
        Returns ExportResult with bytes on success.
        Never raises — errors are captured in ExportResult.error.
        """
        try:
            errors = req.validate()
            if errors:
                return ExportResult(success=False, error=" | ".join(errors))

            theme = get_theme(req.theme)
            prs = self._init_presentation()
            self._build_slides(prs, req, theme)
            data = self._serialize(prs)
            self._validate(data)

            n = len(prs.slides._sldIdLst)
            log.info(f"✅ PPTX built: {n} slides | theme={req.theme} | font={self._font} | {len(data):,} bytes")
            return ExportResult(
                success=True,
                data=data,
                slide_count=n,
                font_used=self._font,
            )

        except Exception as exc:
            log.error(f"Export pipeline error: {exc}", exc_info=True)
            return ExportResult(success=False, error=str(exc))

    # ── Private helpers ───────────────────────────────────────────────

    def _init_presentation(self) -> Presentation:
        prs = Presentation()
        prs.slide_width = Cm(W_CM)
        prs.slide_height = Cm(H_CM)
        return prs

    def _build_slides(self, prs: Presentation, req: PresentationRequest, T: Theme):
        cfg = req.slides

        # Cover is always built
        make_cover(prs, req, T)

        if cfg.intro and (req.intro_overview or req.intro_approach):
            make_intro(prs, req, T)

        chapters = req.chapters
        if cfg.plan and chapters:
            make_plan(prs, req, T)

        if cfg.problem and (req.main_problem or req.main_question or req.sub_questions):
            make_problem(prs, req, T)

        if cfg.objectives and (req.objectives or req.hypotheses):
            make_objectives(prs, req, T)

        if cfg.importance and (req.importance or req.reasons):
            make_importance(prs, req, T)

        if cfg.methodology and (req.methodology or req.sample_type or req.tool):
            make_methodology(prs, req, T)

        if cfg.kpi and req.stats:
            make_stats(prs, req, T)

        if cfg.results and req.main_results:
            make_results(prs, req, T)

        if cfg.conclusion and req.general_conclusion:
            make_conclusion(prs, req, T)

        if cfg.recommendations and req.recommendations:
            make_recommendations(prs, req, T)

        if cfg.future and req.future_work:
            make_future(prs, req, T)

        if cfg.references and req.references:
            make_references(prs, req, T)

        if cfg.thankyou:
            make_final(prs, req, T)

    def _serialize(self, prs: Presentation) -> bytes:
        """Save to in-memory buffer. No temp files."""
        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    def _validate(self, data: bytes):
        """Basic sanity checks on the output bytes."""
        if len(data) < MIN_VALID_SIZE:
            raise ValueError(f"Output too small ({len(data)} bytes) — likely corrupt")
        # PPTX is a ZIP — check magic bytes
        if not data.startswith(b'PK'):
            raise ValueError("Output is not a valid ZIP/PPTX (missing PK header)")


# ── Module-level singleton ────────────────────────────────────────────
_pipeline: PPTXExportPipeline | None = None


def get_pipeline() -> PPTXExportPipeline:
    """Lazy singleton — initialized once on first call."""
    global _pipeline
    if _pipeline is None:
        _pipeline = PPTXExportPipeline()
    return _pipeline
