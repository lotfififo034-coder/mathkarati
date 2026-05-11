"""
Drawing Primitives — مذكرتي Pro v17
Low-level, deterministic shape/text builders.
All functions are pure: they mutate only the slide they're given.
No global state. No side effects outside the slide.
"""
from __future__ import annotations

from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# Slide dimensions (cm) — 16:9 1920×1080
W, H = 33.867, 19.05


# ── Unit helpers ────────────────────────────────────────────────────
def cm(v: float) -> int:
    return int(Cm(v))


def pt(v: float) -> int:
    return int(Pt(v))


# ── Shape builders ──────────────────────────────────────────────────
def rect(slide, x: float, y: float, w: float, h: float,
         fill: RGBColor, line: RGBColor | None = None, line_w: float = 0.5):
    """Solid rectangle. Returns shape or None if degenerate."""
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = pt(line_w)
    else:
        s.line.fill.background()
    return s


def rrect(slide, x: float, y: float, w: float, h: float,
          fill: RGBColor, radius_pct: int = 8,
          line: RGBColor | None = None, line_w: float = 0.5):
    """Rounded rectangle."""
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(5, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = pt(line_w)
    else:
        s.line.fill.background()
    try:
        adj = s.adjustments
        if adj and len(adj) > 0:
            adj[0] = max(0, min(50, radius_pct)) * 1000
    except Exception:
        pass
    return s


def oval(slide, x: float, y: float, w: float, h: float,
         fill: RGBColor, alpha: int = 100):
    """Ellipse/circle with optional transparency."""
    if w <= 0 or h <= 0:
        return None
    s = slide.shapes.add_shape(9, cm(x), cm(y), cm(w), cm(h))
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    if alpha < 100:
        _set_solid_alpha(s, alpha)
    return s


def bg(slide, color: RGBColor):
    """Fill entire slide background."""
    rect(slide, 0, 0, W, H, color)


def hline(slide, x: float, y: float, w: float, color: RGBColor, thickness: float = 0.08):
    rect(slide, x, y, w, thickness, color)


def vline(slide, x: float, y: float, h2: float, color: RGBColor, thickness: float = 0.08):
    rect(slide, x, y, thickness, h2, color)


# ── XML-level fill helpers ───────────────────────────────────────────
def _set_solid_alpha(shape, alpha_pct: int):
    try:
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        fld = spPr.find('.//' + qn('a:solidFill'))
        if fld is not None:
            srgb = fld.find(qn('a:srgbClr'))
            if srgb is not None:
                for e in srgb.findall(qn('a:alpha')):
                    srgb.remove(e)
                alp = etree.SubElement(srgb, qn('a:alpha'))
                alp.set('val', str(int(alpha_pct * 1000)))
    except Exception:
        pass


def gradient_fill(shape, c1: str, c2: str, angle: float = 90):
    """Apply a real linear gradient via XML. c1/c2 are hex strings."""
    try:
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        # Remove any existing fill
        for tag in [qn('a:solidFill'), qn('a:gradFill'), qn('a:noFill'),
                    qn('a:pattFill'), qn('a:blipFill')]:
            for el in spPr.findall(tag):
                spPr.remove(el)
        # Build gradFill
        grad = etree.SubElement(spPr, qn('a:gradFill'))
        gsLst = etree.SubElement(grad, qn('a:gsLst'))
        gs0 = etree.SubElement(gsLst, qn('a:gs'))
        gs0.set('pos', '0')
        sc0 = etree.SubElement(gs0, qn('a:srgbClr'))
        sc0.set('val', c1.lstrip('#'))
        gs1 = etree.SubElement(gsLst, qn('a:gs'))
        gs1.set('pos', '100000')
        sc1 = etree.SubElement(gs1, qn('a:srgbClr'))
        sc1.set('val', c2.lstrip('#'))
        lin = etree.SubElement(grad, qn('a:lin'))
        lin.set('ang', str(int(angle * 60000)))
        lin.set('scaled', '0')
    except Exception:
        pass


def gradient_rect(slide, x: float, y: float, w: float, h: float,
                  c1: str, c2: str, angle: float = 0):
    """Rectangle with a real gradient fill."""
    from pptx.dml.color import RGBColor
    c1h = c1.lstrip('#')
    fill_color = RGBColor(int(c1h[0:2], 16), int(c1h[2:4], 16), int(c1h[4:6], 16))
    s = rect(slide, x, y, w, h, fill_color)
    if s:
        gradient_fill(s, c1, c2, angle)
    return s


def shadow(shape, blur: float = 16, dist: float = 5,
           angle: float = 135, alpha: float = 0.22, color: str = "000000"):
    """Add outer drop shadow via XML."""
    try:
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        # Remove old effect list
        for old in spPr.findall(qn('a:effectLst')):
            spPr.remove(old)
        eLst = etree.SubElement(spPr, qn('a:effectLst'))
        shdw = etree.SubElement(eLst, qn('a:outerShdw'))
        shdw.set('blurRad', str(int(blur * 12700)))
        shdw.set('dist', str(int(dist * 12700)))
        shdw.set('dir', str(int(angle * 60000)))
        shdw.set('algn', 'tl')
        srgb = etree.SubElement(shdw, qn('a:srgbClr'))
        srgb.set('val', color)
        alp = etree.SubElement(srgb, qn('a:alpha'))
        alp.set('val', str(int(alpha * 100000)))
    except Exception:
        pass


def glow(shape, color: str, radius: float = 8, alpha: float = 0.40):
    """Add glow effect via XML."""
    try:
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        eLst = spPr.find(qn('a:effectLst'))
        if eLst is None:
            eLst = etree.SubElement(spPr, qn('a:effectLst'))
        g = etree.SubElement(eLst, qn('a:glow'))
        g.set('rad', str(int(radius * 12700)))
        srgb = etree.SubElement(g, qn('a:srgbClr'))
        srgb.set('val', color.lstrip('#'))
        alp = etree.SubElement(srgb, qn('a:alpha'))
        alp.set('val', str(int(alpha * 100000)))
    except Exception:
        pass


# ── Text ─────────────────────────────────────────────────────────────
def txt(slide, text: str,
        x: float, y: float, w: float, h: float,
        font: str = "Cairo",
        size: float = 14,
        bold: bool = False,
        italic: bool = False,
        color: RGBColor | None = None,
        align=PP_ALIGN.RIGHT,
        margin: float = 0.12,
        rtl: bool = True,
        spacing: float | None = None):
    """Add a textbox. Returns None if text/dimensions are invalid."""
    if not text or w <= 0 or h <= 0:
        return None
    tb = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = cm(margin)
    tf.margin_right = cm(margin)
    tf.margin_top = cm(0.04)
    tf.margin_bottom = cm(0.04)
    p = tf.paragraphs[0]
    p.alignment = align
    if spacing is not None:
        try:
            p.line_spacing = Pt(spacing)
        except Exception:
            pass
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return tb


def txt_multiline(slide, lines: list[str],
                  x: float, y: float, w: float, h: float,
                  font: str = "Cairo",
                  size: float = 12,
                  bold: bool = False,
                  color: RGBColor | None = None,
                  align=PP_ALIGN.RIGHT,
                  margin: float = 0.12,
                  line_spacing: float = 1.4):
    """Add a textbox with multiple paragraphs."""
    if not lines or w <= 0 or h <= 0:
        return None
    tb = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = cm(margin)
    tf.margin_right = cm(margin)
    tf.margin_top = cm(0.04)
    tf.margin_bottom = cm(0.04)

    for i, line in enumerate(lines):
        if not line:
            continue
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        try:
            p.line_spacing = Pt(size * line_spacing)
        except Exception:
            pass
        run = p.add_run()
        run.text = str(line)
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        if color:
            run.font.color.rgb = color
    return tb


def blank_slide(prs):
    """Add a completely blank slide (layout 6 = blank)."""
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)
