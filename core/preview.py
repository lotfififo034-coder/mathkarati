"""
core/preview.py — توليد صور معاينة حقيقية من PPTX مع علامة مائية
يستخدم LibreOffice لتحويل دقيق، مع Pillow كـ fallback.
"""
import base64
import io
import logging
import os
import shutil
import subprocess
import tempfile
import threading
from pathlib import Path
from typing import List, Optional

log = logging.getLogger(__name__)

# ── Cache ─────────────────────────────────────────────────────────────────────
_cache: dict = {}
_cache_lock = threading.Lock()
MAX_PREVIEW_SLIDES = 3   # نعرض أول 3 شرائح فقط كمعاينة


def get_cached_preview(presentation_id: str) -> Optional[List[str]]:
    with _cache_lock:
        return _cache.get(presentation_id)


def set_cached_preview(presentation_id: str, slides: List[str]) -> None:
    with _cache_lock:
        _cache[presentation_id] = slides


# ── Main entry ────────────────────────────────────────────────────────────────
def pptx_to_preview_images(pptx_path: str, watermark: bool = True) -> List[str]:
    """
    يحوّل أول MAX_PREVIEW_SLIDES شرائح من ملف PPTX إلى JPEG base64.
    يجرّب LibreOffice أولاً، ثم Pillow كـ fallback.
    """
    try:
        slides = _render_libreoffice(pptx_path)
        if slides:
            log.info(f"LibreOffice rendered {len(slides)} slides")
            if watermark:
                slides = [_add_watermark(s) for s in slides]
            return slides
    except Exception as exc:
        log.warning(f"LibreOffice render failed: {exc}, trying Pillow fallback")

    try:
        slides = _render_pillow_fallback(pptx_path)
        if watermark:
            slides = [_add_watermark(s) for s in slides]
        return slides
    except Exception as exc:
        log.warning(f"Pillow fallback also failed: {exc}")
        return []


# ── LibreOffice renderer ──────────────────────────────────────────────────────
def _find_libreoffice() -> Optional[str]:
    for cmd in ("libreoffice", "libreoffice7.6", "libreoffice7.5",
                "libreoffice7.4", "soffice"):
        path = shutil.which(cmd)
        if path:
            return path
    # مسارات شائعة
    for p in ("/usr/bin/libreoffice", "/usr/lib/libreoffice/program/soffice",
              "/opt/libreoffice/program/soffice"):
        if os.path.exists(p):
            return p
    return None


def _render_libreoffice(pptx_path: str) -> List[str]:
    lo = _find_libreoffice()
    if not lo:
        raise RuntimeError("LibreOffice not found")

    with tempfile.TemporaryDirectory() as tmpdir:
        # نسخ الملف للـ tmpdir لتجنب مشاكل الصلاحيات
        src = Path(tmpdir) / "presentation.pptx"
        shutil.copy2(pptx_path, src)

        # تحويل إلى PNG
        result = subprocess.run(
            [lo, "--headless", "--norestore", "--convert-to", "png",
             "--outdir", tmpdir, str(src)],
            capture_output=True, text=True, timeout=120,
            env={**os.environ, "HOME": tmpdir, "TMPDIR": tmpdir}
        )
        log.info(f"LO stdout: {result.stdout[:200]}")
        if result.returncode != 0:
            log.warning(f"LO stderr: {result.stderr[:300]}")

        # البحث عن ملفات PNG المُنتجة
        # LibreOffice تنتج: presentation.png (صفحة واحدة) أو presentation1.png, presentation2.png...
        pngs = sorted([
            f for f in Path(tmpdir).glob("presentation*.png")
        ], key=lambda x: x.name)

        if not pngs:
            raise RuntimeError("LibreOffice produced no PNG files")

        slides = []
        for png in pngs[:MAX_PREVIEW_SLIDES]:
            from PIL import Image
            img = Image.open(png).convert("RGB")
            # resize لتوفير الحجم
            img.thumbnail((1280, 720), Image.LANCZOS)
            buf = io.BytesIO()
            img.save(buf, format="JPEG", quality=82, optimize=True)
            slides.append(base64.b64encode(buf.getvalue()).decode("ascii"))

        return slides


# ── Pillow fallback ───────────────────────────────────────────────────────────
def _render_pillow_fallback(pptx_path: str) -> List[str]:
    """
    تحويل تقريبي بـ Pillow: يرسم خلفية + نصوص الشرائح.
    ليس دقيقاً لكن أفضل من لا شيء.
    """
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont

    prs = Presentation(pptx_path)
    EMU = 914400
    DPI = 96
    W = max(960, int(prs.slide_width / EMU * DPI))
    H = max(540, int(prs.slide_height / EMU * DPI))

    results = []
    for slide in list(prs.slides)[:MAX_PREVIEW_SLIDES]:
        img = _draw_slide(slide, W, H)
        buf = io.BytesIO()
        img.save(buf, format="JPEG", quality=80)
        results.append(base64.b64encode(buf.getvalue()).decode("ascii"))
    return results


def _draw_slide(slide, W: int, H: int):
    from PIL import Image, ImageDraw, ImageFont

    bg = _slide_bg_color(slide)
    img = Image.new("RGB", (W, H), bg)
    draw = ImageDraw.Draw(img)

    for shape in slide.shapes:
        try:
            if not shape.has_text_frame:
                continue
            left = int((shape.left or 0) / 914400 * 96)
            top  = int((shape.top  or 0) / 914400 * 96)
            y = top
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    y += 12
                    continue
                fsize = 16
                color = (30, 30, 30)
                for run in para.runs:
                    try:
                        if run.font.size: fsize = max(8, min(int(run.font.size/12700), 60))
                        if run.font.color.rgb:
                            c = run.font.color.rgb
                            color = (c.r, c.g, c.b)
                    except Exception:
                        pass
                try:
                    font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", fsize)
                except Exception:
                    font = ImageFont.load_default()
                draw.text((left + 6, y), text, fill=color, font=font)
                y += fsize + 5
        except Exception:
            continue
    return img


def _slide_bg_color(slide) -> tuple:
    try:
        fill = slide.background.fill
        if fill.type is not None:
            rgb = fill.fore_color.rgb
            return (rgb.r, rgb.g, rgb.b)
    except Exception:
        pass
    return (255, 255, 255)


# ── Watermark ─────────────────────────────────────────────────────────────────
def _add_watermark(b64_jpeg: str) -> str:
    """يضيف علامة مائية خفيفة على صورة JPEG مُرمَّزة بـ base64."""
    from PIL import Image, ImageDraw, ImageFont

    data = base64.b64decode(b64_jpeg)
    img = Image.open(io.BytesIO(data)).convert("RGBA")
    W, H = img.size

    overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)

    font_size = max(22, W // 28)
    font = None
    for fp in [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
    ]:
        if os.path.exists(fp):
            try:
                font = ImageFont.truetype(fp, font_size)
                break
            except Exception:
                pass
    if font is None:
        font = ImageFont.load_default()

    text = "مذكرتي Pro — معاينة فقط"

    # علامة مائية قطرية في عدة مواضع
    import math
    positions = [
        (W * 0.15, H * 0.25),
        (W * 0.45, H * 0.50),
        (W * 0.20, H * 0.70),
        (W * 0.55, H * 0.20),
        (W * 0.35, H * 0.80),
    ]
    # رسم مؤقت لقياس النص
    tmp_draw = ImageDraw.Draw(Image.new("RGBA", (1, 1)))
    try:
        bbox = tmp_draw.textbbox((0, 0), text, font=font)
        tw = bbox[2] - bbox[0]
        th = bbox[3] - bbox[1]
    except Exception:
        tw, th = font_size * len(text) // 2, font_size

    # إنشاء طبقة دوارة للعلامة
    txt_img = Image.new("RGBA", (tw + 20, th + 20), (0, 0, 0, 0))
    txt_draw = ImageDraw.Draw(txt_img)
    txt_draw.text((10, 10), text, fill=(255, 255, 255, 70), font=font)
    rotated = txt_img.rotate(-30, expand=True)

    for (px, py) in positions:
        x = int(px - rotated.width // 2)
        y = int(py - rotated.height // 2)
        overlay.paste(rotated, (x, y), rotated)

    combined = Image.alpha_composite(img, overlay).convert("RGB")
    buf = io.BytesIO()
    combined.save(buf, format="JPEG", quality=82, optimize=True)
    return base64.b64encode(buf.getvalue()).decode("ascii")
