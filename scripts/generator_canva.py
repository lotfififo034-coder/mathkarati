"""
مذكرتي Pro — Classic Engine CANVA LEVEL
========================================
تصاميم سينمائية احترافية مستوحاة من Canva Pro + Behance Portfolio
3 عائلات بصرية × 8 لوحات ألوان = 24 إصدار فريد

معمارية التصميم:
- NOIR    : فاخر داكن — Palatino + Cairo — ذهبي رصاصي
- VIVID   : حيوي مضيء — Trebuchet + Cairo — ألوان بارزة
- MINIMAL : نظيف أبيض — Calibri + Cairo — خطوط رفيعة
"""

import sys, json, math, datetime
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ─── Dimensions (widescreen 33.87 × 19.05 cm) ─────────────────────
W, H   = 33.87 / 2.54, 19.05 / 2.54  # in inches for Cm()
W, H   = 13.33, 7.50
MX, MY = 0.65, 0.50

# ─── Helpers ──────────────────────────────────────────────────────
def rgb(r, g, b): return RGBColor(r, g, b)
def hx(h):        return RGBColor.from_string(h)
def safe(v, fb=""): return str(v).strip() if v else fb
def clamp(v,lo,hi): return max(lo,min(hi,v))
def blank(prs):    return prs.slides.add_slide(prs.slide_layouts[6])

def cm(v): return Cm(v)

def rect(slide, x,y,w,h, fill, alpha=1.0, line_color=None, line_w=0.5, radius=0):
    """رسم مستطيل مع دعم الحواف الدائرية وخيوط الحدود"""
    if w<=0 or h<=0: return
    shape = slide.shapes.add_shape(1, cm(x),cm(y),cm(w),cm(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if alpha < 1.0:
        shape.fill.fore_color.theme_color = None
        t = int((1-alpha)*100000)
        shape.fill.fore_color._xClr.attrib['lumMod'] = str(t)
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def rrect(slide, x,y,w,h, fill, radius_pct=8, alpha=1.0, line_color=None):
    """مستطيل بزوايا دائرية"""
    if w<=0 or h<=0: return
    shape = slide.shapes.add_shape(5, cm(x),cm(y),cm(w),cm(h))  # 5=roundedRectangle
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    if line_color:
        shape.line.color.rgb = line_color; shape.line.width = Pt(0.6)
    else:
        shape.line.fill.background()
    # Set corner radius
    adj = shape.adjustments
    if adj: adj[0] = int(radius_pct * 1000)
    return shape

def oval(slide, x,y,w,h, fill, alpha=1.0):
    if w<=0 or h<=0: return
    shape = slide.shapes.add_shape(9, cm(x),cm(y),cm(w),cm(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def line_h(slide, x,y,w, color, h=0.04):
    rect(slide, x,y,w,h, color)

def line_v(slide, x,y,h_cm, color, w=0.04):
    rect(slide, x,y,w,h_cm, color)

def bg(slide, color):
    rect(slide, 0,0,W,H, color)

def txt(slide, text, x,y,w,h,
        font="Cairo", size=13, bold=False, italic=False,
        color=None, align=PP_ALIGN.RIGHT, mg=0.08,
        rtl=True, valign="top", line_spacing=None):
    if w<=0 or h<=0 or not text: return None
    tb = slide.shapes.add_textbox(cm(x),cm(y),cm(w),cm(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left=cm(mg); tf.margin_right=cm(mg)
    tf.margin_top=cm(0.03); tf.margin_bottom=cm(0.03)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color: run.font.color.rgb = color
    return tb

def shadow(shape, blur=8, dist=3, angle=135, color="000000", alpha=0.15):
    """إضافة ظل للشكل"""
    try:
        sp = shape._element
        spPr = sp.find(qn('p:spPr'))
        if spPr is None: return
        effectLst = etree.SubElement(spPr, qn('a:effectLst'))
        outerShdw = etree.SubElement(effectLst, qn('a:outerShdw'))
        outerShdw.set('blurRad', str(blur*12700))
        outerShdw.set('dist', str(dist*12700))
        outerShdw.set('dir', str(int(angle*60000)))
        srgbClr = etree.SubElement(outerShdw, qn('a:srgbClr'))
        srgbClr.set('val', color)
        alpha_elem = etree.SubElement(srgbClr, qn('a:alpha'))
        alpha_elem.set('val', str(int(alpha*100000)))
    except: pass

# ═══════════════════════════════════════════════════════════════════
# PALETTES — 8 لوحات ألوان
# ═══════════════════════════════════════════════════════════════════
def P(bg_dark, bg_mid, bg_light, accent, accent2, txt_light, txt_dark, txt_mid,
      chart_colors, hf, bf, layout_family):
    return {
        "D":  hx(bg_dark),   "M":  hx(bg_mid),  "L":  hx(bg_light),
        "A":  hx(accent),    "A2": hx(accent2),
        "TL": hx(txt_light), "TD": hx(txt_dark), "TM": hx(txt_mid),
        "CB": hx("FFFFFF"),  "CE": hx("E8EDF2"),
        "SC": [hx(c) for c in chart_colors],
        "HF": hf, "BF": bf, "FAM": layout_family,
    }

PALETTES = {
# ── NOIR family (داكن فاخر) ────────────────────────────────────────
"navy_gold": P(
    "07172F","0E274D","F4F7FB","C6A03C","E8C97B","FFFFFF","07172F","94A3B8",
    ["C6A03C","4AB3FF","E8C97B","1A4A8A","FF6B6B","00D4AA","F59E0B"],
    "Palatino Linotype","Cairo","NOIR"),

"midnight_purple": P(
    "12052E","2D106B","F5F3FF","C07AFF","DDB3FF","FFFFFF","12052E","9CA3AF",
    ["C07AFF","FF6B6B","DDB3FF","7B3FE0","00D4AA","F59E0B","4AB3FF"],
    "Georgia","Cairo","NOIR"),

"forest": P(
    "0F2D1E","1E4D36","F0FDF4","86BB56","B4D984","FFFFFF","0F2D1E","6B7A6B",
    ["86BB56","4AB3FF","B4D984","1E4D36","FF6B6B","F59E0B","00D4AA"],
    "Georgia","Cairo","NOIR"),

"sand_gold": P(
    "2C1A0E","5C3D20","FFFBF5","C8861F","E8B964","FFFFFF","2C1A0E","92724A",
    ["C8861F","4AB3FF","E8B964","5C3D20","FF6B6B","00D4AA","86BB56"],
    "Palatino Linotype","Cairo","NOIR"),

# ── VIVID family (حيوي مضيء) ──────────────────────────────────────
"dark_teal": P(
    "0B1F2E","103A50","EBF8FF","00C9A7","67E8D3","FFFFFF","061A28","5FA8B8",
    ["00C9A7","FF6B35","67E8D3","0993C3","C07AFF","F59E0B","86BB56"],
    "Trebuchet MS","Cairo","VIVID"),

"charcoal_orange": P(
    "1A1A2E","2D2D44","FFF8F5","FF6B35","FFA070","FFFFFF","1A1A2E","7A7A96",
    ["FF6B35","00C9A7","FFA070","2D2D44","C07AFF","4AB3FF","86BB56"],
    "Trebuchet MS","Cairo","VIVID"),

"burgundy": P(
    "3A0018","6B1537","FFF5F7","E84393","FFAACC","FFFFFF","3A0018","9A5070",
    ["E84393","4AB3FF","FFAACC","6B1537","00C9A7","F59E0B","C07AFF"],
    "Georgia","Cairo","VIVID"),

# ── MINIMAL family (نظيف أبيض) ────────────────────────────────────
"ice_blue": P(
    "0A254A","1A4A8A","F0F7FF","0066CC","4AB3FF","FFFFFF","0A254A","4A6A8A",
    ["0066CC","FF6B35","4AB3FF","1A4A8A","00C9A7","C07AFF","86BB56"],
    "Calibri","Cairo","MINIMAL"),
}


# ═══════════════════════════════════════════════════════════════════
# DESIGN SYSTEM — functions shared across all slides
# ═══════════════════════════════════════════════════════════════════

def deco_blob(slide, T, cx, cy, r, alpha=0.06):
    """دائرة ديكورية خفية"""
    oval(slide, cx-r, cy-r, r*2, r*2, T["M"])

def stripe(slide, x,y,w,h, color, n=6, spacing=0.28):
    """خطوط مائلة ديكورية (Canva style)"""
    for i in range(n):
        xi = x + i*spacing
        if xi < x+w:
            rect(slide, xi, y, 0.06, h, color)

def header_bar(slide, T, title_ar, sub_en="", dark=True, h=1.90):
    """شريط عنوان موحد مع overline + خط فاصل"""
    bar_c = T["D"] if dark else T["L"]
    bg(slide, bar_c)
    # شريط accent رفيع أعلى
    rect(slide, 0, 0, W, 0.10, T["A"])
    # overline
    txt(slide, sub_en.upper() if sub_en else "", MX, 0.18, W-MX*2, 0.32,
        font="Calibri", size=8.5, bold=True, color=T["A"],
        align=PP_ALIGN.LEFT, rtl=False)
    # العنوان الرئيسي
    c = T["TL"] if dark else T["TD"]
    txt(slide, title_ar, MX, 0.42, W-MX*2, 1.20,
        font=T["HF"], size=30, bold=True, color=c,
        align=PP_ALIGN.RIGHT, rtl=True)
    # خط فاصل gradient (نحاكيه بمستطيلات)
    line_h(slide, MX, h-0.12, 1.80, T["A"], 0.06)
    line_h(slide, MX+1.80, h-0.12, W-MX*2-1.80, T["A"] if dark else T["CE"], 0.06)
    return h  # content starts here

def pill(slide, x,y,w,h, text, bg_c, text_c, font="Cairo", size=10, bold=True):
    """بادج/pill مستطيل بزوايا دائرية"""
    s = rrect(slide, x,y,w,h, bg_c, radius_pct=50)
    txt(slide, text, x,y,w,h, font=font, size=size, bold=bold,
        color=text_c, align=PP_ALIGN.CENTER, mg=0.05)

def number_badge(slide, x,y,size_cm, n, bg_c, text_c):
    """دائرة رقم (01 02 03)"""
    oval(slide, x,y,size_cm,size_cm, bg_c)
    txt(slide, "%02d" % n, x,y,size_cm,size_cm,
        font="Calibri", size=int(size_cm*10), bold=True,
        color=text_c, align=PP_ALIGN.CENTER)

def card_dark(slide, x,y,w,h, T, accent_color=None, show_shadow=True):
    """بطاقة داكنة مع شريط accent"""
    s = rrect(slide, x,y,w,h, T["M"], radius_pct=6)
    if show_shadow: shadow(s)
    if accent_color:
        rect(slide, x, y, w, 0.08, accent_color)
    return s

def card_light(slide, x,y,w,h, T, accent_color=None):
    """بطاقة فاتحة بيضاء مع حدود"""
    s = rrect(slide, x,y,w,h, T["CB"], radius_pct=6, line_color=T["CE"])
    if s: shadow(s, blur=5, dist=2, alpha=0.08)
    if accent_color:
        rect(slide, x, y, 0.08, h, accent_color)
    return s

def kpi_card(slide, x,y,w,h, T, value, label, color, dark=True):
    """بطاقة KPI سينمائية"""
    bg_c = T["M"] if dark else T["CB"]
    s = rrect(slide, x,y,w,h, bg_c, radius_pct=8)
    if s: shadow(s, blur=10, dist=3, alpha=0.14)
    # شريط علوي ملون
    rect(slide, x,y,w,0.10, color)
    # القيمة
    v = safe(value)
    vsize = clamp(52 - max(0,len(v)-4)*7, 28, 52)
    txt(slide, v, x+0.10, y+0.22, w-0.20, h*0.52,
        font="Calibri", size=vsize, bold=True,
        color=color, align=PP_ALIGN.CENTER)
    # خط فاصل
    line_h(slide, x+0.20, y+h*0.66, w-0.40, color, 0.025)
    # التسمية
    tc = T["TL"] if dark else T["TD"]
    txt(slide, safe(label), x+0.10, y+h*0.70, w-0.20, h*0.26,
        font=T["BF"], size=11, color=tc, align=PP_ALIGN.CENTER)

def bullet_row(slide, x,y,w,h, T, n, text, color, alt=False, dark=True):
    """صف قائمة مرقم Canva-style"""
    bg_c = T["M"] if dark else (T["L"] if alt else T["CB"])
    s = rrect(slide, x,y,w,h, bg_c, radius_pct=5)
    # شريط يمين
    rect(slide, x, y, 0.10, h, color)
    # رقم
    oval(slide, x+0.20, y+(h-0.55)/2, 0.55, 0.55, color)
    txt(slide, str(n), x+0.20, y+(h-0.55)/2, 0.55, 0.55,
        font="Calibri", size=13, bold=True,
        color=T["D"], align=PP_ALIGN.CENTER)
    # نص
    tc = T["TL"] if dark else T["TD"]
    txt(slide, safe(text), x+0.90, y+0.08, w-1.02, h-0.16,
        font=T["BF"], size=12, color=tc,
        align=PP_ALIGN.RIGHT, rtl=True)

def chapter_card(slide, x,y,w,h, T, num, title, sections, color):
    """بطاقة فصل في خطة الدراسة"""
    s = rrect(slide, x,y,w,h, T["M"], radius_pct=7)
    if s: shadow(s)
    # رأس البطاقة
    rect(slide, x,y,w,0.08, color)
    # رقم الفصل
    txt(slide, "F%d" % num, x+0.18, y+0.12, w-0.30, 0.60,
        font="Calibri", size=22, bold=True,
        color=color, align=PP_ALIGN.RIGHT)
    line_h(slide, x+0.18, y+0.76, w-0.30, color, 0.025)
    # عنوان الفصل
    txt(slide, safe(title), x+0.18, y+0.84, w-0.30,
        0.92 if sections else h-0.96,
        font=T["BF"], size=12, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    # المباحث
    if sections:
        secs = [s for s in sections if s][:5]
        sh = (h - 1.82) / max(len(secs),1)
        for i, sec in enumerate(secs):
            sy = y + 1.82 + i*sh
            oval(slide, x+0.24, sy+sh*0.35, 0.10, 0.10, color)
            txt(slide, safe(sec), x+0.42, sy+0.04, w-0.56, sh-0.08,
                font=T["BF"], size=10, color=T["TM"],
                align=PP_ALIGN.RIGHT, rtl=True)


# ═══════════════════════════════════════════════════════════════════
# COVER FAMILIES
# ═══════════════════════════════════════════════════════════════════

def cover_noir(slide, T, data):
    """غلاف NOIR — شريط جانبي داكن + هندسة دائرية"""
    bg(slide, T["D"])
    # خلفية هندسية
    oval(slide, W*0.55, -3.0, W*0.70, H*1.05, T["M"])
    oval(slide, W*0.72, -1.0, W*0.42, H*0.80, T["A"])
    deco_blob(slide, T, -2.0, H*0.70, 9.0, 0.40)

    # شريط يميني رأسي
    rect(slide, W-0.80, 0, 0.80, H, T["A"])
    # شريط سفلي
    rect(slide, 0, H-0.60, W-0.80, 0.60, T["A"])

    # نطاق جامعة
    rect(slide, 0, 0, W-0.80, 2.70, T["M"])
    line_h(slide, 0, 2.70, W-0.80, T["A"], 0.06)

    # رقم السنة ضخم ديكوري
    yr = safe(data.get("year","")).split("–")[-1].split("-")[-1]
    if yr:
        txt(slide, yr, 0.10, H*0.22, W*0.55, H*0.60,
            font="Calibri", size=180, bold=True,
            color=T["M"], align=PP_ALIGN.LEFT, rtl=False)

    # ── محتوى نصي ────────────────────────────────────
    # اسم الجامعة
    txt(slide, safe(data.get("university","")), 0.30, 0.18, W-1.40, 0.88,
        font=T["BF"], size=14, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    # الكلية والقسم
    fac = " · ".join(filter(None,[safe(data.get("faculty","")),safe(data.get("department",""))]))
    if fac:
        txt(slide, fac, 0.30, 1.10, W-1.40, 0.56,
            font=T["BF"], size=11, color=T["A"],
            align=PP_ALIGN.RIGHT, rtl=True)
    # المستوى
    pill(slide, 0.30, 1.76, 2.80, 0.56,
         "مذكرة تخرج  ·  " + safe(data.get("level","ماستر 2")),
         T["A"], T["D"], size=11)

    # عنوان المذكرة
    txt(slide, safe(data.get("titleAr","")), 0.30, 3.10, W-1.40, 3.30,
        font=T["BF"], size=22, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)

    # العنوان الفرنسي
    if data.get("titleFr"):
        line_h(slide, 0.30, 6.52, W-1.40, T["A"], 0.04)
        txt(slide, safe(data.get("titleFr","")), 0.30, 6.62, W-1.40, 0.56,
            font="Calibri", size=11, italic=True,
            color=T["A"], align=PP_ALIGN.LEFT, rtl=False)

    # باحث + مشرف
    line_h(slide, 0, H-0.60, W-0.80, T["A"], 0.06)
    hw = (W-1.40)/2
    txt(slide, "إعداد", 0.30, H-0.54, hw*0.40, 0.42,
        font=T["BF"], size=9, color=T["A"], align=PP_ALIGN.RIGHT)
    txt(slide, safe(data.get("studentName","")), 0.30, H-0.54, hw-0.10, 0.42,
        font=T["BF"], size=13, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    txt(slide, "إشراف", hw+0.50, H-0.54, hw*0.40, 0.42,
        font=T["BF"], size=9, color=T["A"], align=PP_ALIGN.RIGHT)
    txt(slide, safe(data.get("supervisor","")), hw+0.50, H-0.54, hw-0.20, 0.42,
        font=T["BF"], size=13, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)

    # الكلمات المفتاحية + التاريخ
    kw = safe(data.get("keywords",""))
    if kw:
        rect(slide, 0, H-0.60, W-0.80, 0.60, T["A"])
        txt(slide, "🔑 " + kw, 0.30, H-0.56, W-1.40, 0.48,
            font=T["BF"], size=9, italic=True,
            color=T["D"], align=PP_ALIGN.RIGHT)


def cover_vivid(slide, T, data):
    """غلاف VIVID — خلفية مقسمة + هندسة جرئية"""
    bg(slide, T["D"])
    # لوحة يمين فاتحة
    split_x = W * 0.42
    rect(slide, split_x, 0, W-split_x, H, T["M"])
    # مثلثات ديكورية
    oval(slide, split_x-3, H*0.10, 6.0, 6.0, T["A"])
    oval(slide, split_x+1, H*0.55, 4.0, 4.0, T["D"])
    oval(slide, W-4, -2, 7, 7, T["A"])

    # شريط accent رأسي
    rect(slide, split_x-0.10, 0, 0.20, H, T["A"])
    # شريط accent أفقي علوي
    rect(slide, 0, 0, W, 0.12, T["A"])

    # محتوى اليسار (داكن)
    txt(slide, safe(data.get("university","")), 0.40, 0.26, split_x-0.60, 1.60,
        font=T["BF"], size=15, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    fac = " · ".join(filter(None,[safe(data.get("faculty","")),safe(data.get("department",""))]))
    if fac:
        txt(slide, fac, 0.40, 1.94, split_x-0.60, 0.60,
            font=T["BF"], size=11, color=T["A"],
            align=PP_ALIGN.RIGHT, rtl=True)
    pill(slide, 0.40, 2.68, 3.0, 0.56,
         safe(data.get("level","ماستر 2")),
         T["A"], T["D"], size=12, bold=True)

    txt(slide, safe(data.get("studentName","")), 0.40, H-2.0, split_x-0.60, 0.72,
        font=T["BF"], size=18, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    txt(slide, "إشراف: " + safe(data.get("supervisor","")), 0.40, H-1.20, split_x-0.60, 0.52,
        font=T["BF"], size=11, color=T["A"],
        align=PP_ALIGN.RIGHT, rtl=True)
    if data.get("year"):
        pill(slide, 0.40, H-0.60, 1.80, 0.44,
             safe(data.get("year","")), T["M"], T["A"], size=10)

    # محتوى اليمين (عنوان المذكرة)
    rx, rw = split_x+0.40, W-split_x-0.60
    txt(slide, "عنوان المذكرة", rx, 0.30, rw, 0.42,
        font="Calibri", size=9, bold=True,
        color=T["A"], align=PP_ALIGN.LEFT, rtl=False)
    line_h(slide, rx, 0.76, rw, T["A"], 0.06)
    txt(slide, safe(data.get("titleAr","")), rx, 0.90, rw, 4.20,
        font=T["BF"], size=22, bold=True,
        color=T["TD"], align=PP_ALIGN.RIGHT, rtl=True)
    if data.get("titleFr"):
        txt(slide, safe(data.get("titleFr","")), rx, 5.22, rw, 0.70,
            font="Calibri", size=11, italic=True,
            color=T["D"], align=PP_ALIGN.LEFT, rtl=False)
    kw = safe(data.get("keywords",""))
    if kw:
        line_h(slide, rx, H-0.90, rw, T["D"], 0.04)
        txt(slide, kw, rx, H-0.82, rw, 0.60,
            font="Calibri", size=9, italic=True,
            color=T["D"], align=PP_ALIGN.LEFT, rtl=False)


def cover_minimal(slide, T, data):
    """غلاف MINIMAL — أبيض نظيف + تايبوغرافي قوي"""
    bg(slide, T["CB"])
    # شريط جانبي أيسر
    rect(slide, 0, 0, 0.55, H, T["D"])
    rect(slide, 0.55, 0, 0.10, H, T["A"])
    # شريط سفلي
    rect(slide, 0, H-0.12, W, 0.12, T["A"])
    # ظل خفي للشريط
    rect(slide, 0.65, 0.20, 0.04, H-0.32, T["CE"])

    # نص رأسي في الشريط
    txt(slide, (safe(data.get("year","")) + "  ·  مذكرة تخرج").upper(), 0.08, 1.0, H-2.0, 0.44,
        font="Calibri", size=8, bold=True,
        color=T["TL"], align=PP_ALIGN.CENTER, rtl=False)

    # الجامعة
    txt(slide, safe(data.get("university","")), 0.90, 0.26, W-1.20, 1.10,
        font=T["BF"], size=14, bold=True,
        color=T["D"], align=PP_ALIGN.RIGHT, rtl=True)
    fac = " · ".join(filter(None,[safe(data.get("faculty","")),safe(data.get("department",""))]))
    if fac:
        txt(slide, fac, 0.90, 1.40, W-1.20, 0.55,
            font=T["BF"], size=11, color=T["A"],
            align=PP_ALIGN.RIGHT, rtl=True)

    # خط فاصل accent
    line_h(slide, 0.90, 2.08, 3.60, T["A"], 0.10)

    # العنوان — تايبوغرافي ضخم
    txt(slide, safe(data.get("titleAr","")), 0.90, 2.28, W-1.20, 3.20,
        font=T["BF"], size=26, bold=True,
        color=T["D"], align=PP_ALIGN.RIGHT, rtl=True)

    if data.get("titleFr"):
        txt(slide, safe(data.get("titleFr","")), 0.90, 5.58, W-1.20, 0.64,
            font="Calibri", size=12, italic=True,
            color=T["A"], align=PP_ALIGN.LEFT, rtl=False)

    # شريط معلومات أسفل
    rect(slide, 0.65, H-1.20, W-0.65, 1.08, T["D"])
    rect(slide, 0.65, H-1.20, W-0.65, 0.08, T["A"])
    hw = (W-1.20)/2
    txt(slide, "إعداد: "+safe(data.get("studentName","")), 0.90, H-1.06, hw, 0.44,
        font=T["BF"], size=12, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    txt(slide, "إشراف: "+safe(data.get("supervisor","")), 0.90+hw+0.10, H-1.06, hw-0.30, 0.44,
        font=T["BF"], size=12, bold=True,
        color=T["TL"], align=PP_ALIGN.RIGHT, rtl=True)
    pill(slide, W-2.40, H-0.58, 2.0, 0.38,
         safe(data.get("level","ماستر 2")),
         T["A"], T["D"], size=10)


def make_cover(prs, data, T):
    slide = blank(prs)
    fam = T["FAM"]
    if fam == "NOIR":    cover_noir(slide, T, data)
    elif fam == "VIVID": cover_vivid(slide, T, data)
    else:                cover_minimal(slide, T, data)
    return slide


# ═══════════════════════════════════════════════════════════════════
# INTRO SLIDE — المقدمة
# ═══════════════════════════════════════════════════════════════════
def make_intro(prs, data, T):
    slide = blank(prs)
    cy0   = header_bar(slide, T, "المقدمة", "INTRODUCTION", dark=True)
    # ديكور
    oval(slide, W*0.68, H*0.18, W*0.52, H*0.78, T["M"])
    oval(slide, -3, H*0.55, 7, 7, T["A"])

    overview = safe(data.get("introOverview",""))
    approach = safe(data.get("introApproach",""))
    has_ap   = bool(approach)
    oh       = H - cy0 - (2.10 if has_ap else 0.40) - 0.22

    if overview:
        s = rrect(slide, MX, cy0+0.14, W-MX*2, oh, T["M"], radius_pct=7)
        if s: shadow(s)
        rect(slide, MX, cy0+0.14, 0.10, oh, T["A"])
        txt(slide, "❝", MX+0.22, cy0+0.14, 1.40, 0.90,
            font="Georgia", size=48, bold=True, color=T["A"], align=PP_ALIGN.RIGHT)
        txt(slide, overview, MX+0.22, cy0+0.94, W-MX*2-0.34, oh-1.06,
            font=T["BF"], size=13, color=T["TL"],
            align=PP_ALIGN.RIGHT, rtl=True)

    if has_ap:
        ay = H - 1.96
        s2 = rrect(slide, MX, ay, W-MX*2, 1.80, T["A"], radius_pct=7)
        rect(slide, MX, ay, 0.10, 1.80, T["D"])
        txt(slide, "المقاربة النظرية", MX+0.22, ay+0.14, W-MX*2-0.34, 0.48,
            font=T["BF"], size=12, bold=True,
            color=T["D"], align=PP_ALIGN.RIGHT, rtl=True)
        txt(slide, approach, MX+0.22, ay+0.66, W-MX*2-0.34, 1.00,
            font=T["BF"], size=12, color=T["D"],
            align=PP_ALIGN.RIGHT, rtl=True)
    return slide


# ═══════════════════════════════════════════════════════════════════
# PLAN SLIDE — خطة الدراسة
# ═══════════════════════════════════════════════════════════════════
def make_plan(prs, data, T, chapters_data):
    slide = blank(prs)
    cy0   = header_bar(slide, T, "خطة الدراسة", "PLAN D'ETUDE", dark=True)

    chs  = chapters_data[:4]
    n    = len(chs)
    if n == 0: return slide
    gx   = 0.24
    cw   = (W - MX*2 - gx*(n-1)) / n
    cw   = min(cw, 7.0)
    ch   = H - cy0 - 0.40

    card_h = H - cy0 - 0.62
    for i, chap in enumerate(chs):
        cx  = MX + i*(cw+gx)
        sc  = T["SC"][i % len(T["SC"])]
        chapter_card(slide, cx, cy0+0.22, cw, card_h, T,
                     i+1, chap.get("title",""), chap.get("sections",[]), sc)
    return slide


# ═══════════════════════════════════════════════════════════════════
# PROBLEM SLIDE — الإشكالية والتساؤلات
# ═══════════════════════════════════════════════════════════════════
def make_problem(prs, data, T):
    slide  = blank(prs)
    cy0    = header_bar(slide, T, "إشكالية البحث والتساؤلات",
                        "RESEARCH PROBLEM & QUESTIONS", dark=True)
    # ديكور
    oval(slide, W-9, -2, 12, 12, T["M"])

    problem = safe(data.get("mainProblem",""))
    main_q  = safe(data.get("mainQuestion",""))
    subs    = [s for s in data.get("subQuestions",[]) if s][:5]

    if subs:
        # تخطيط: يسار (إشكالية) | يمين (تساؤلات)
        lw = W * 0.455 - MX
        rw = W - MX*2 - lw - 0.30

        # لوحة الإشكالية
        qh = H - cy0 - (1.85 if main_q else 0.40) - 0.30
        s  = rrect(slide, MX, cy0+0.20, lw, qh, T["M"], radius_pct=7)
        if s: shadow(s)
        rect(slide, MX, cy0+0.20, 0.10, qh, T["A"])
        txt(slide, "❝", MX+0.20, cy0+0.22, 1.30, 0.80,
            font="Georgia", size=42, bold=True, color=T["A"])
        txt(slide, problem, MX+0.20, cy0+0.92, lw-0.32, qh-1.04,
            font=T["BF"], size=12.5, color=T["TL"],
            align=PP_ALIGN.RIGHT, rtl=True)

        # التساؤل الرئيسي
        if main_q:
            mq_y = cy0 + 0.20 + qh + 0.18
            s2   = rrect(slide, MX, mq_y, lw, H-mq_y-0.22, T["A"], radius_pct=7)
            txt(slide, "؟", MX+0.12, mq_y+0.10, 0.70, H-mq_y-0.42,
                font="Georgia", size=44, bold=True, color=T["D"])
            txt(slide, main_q, MX+0.90, mq_y+0.14, lw-1.04, H-mq_y-0.38,
                font=T["BF"], size=13, bold=True, color=T["D"],
                align=PP_ALIGN.RIGHT, rtl=True)

        # التساؤلات الفرعية
        rx   = MX + lw + 0.30
        avail = H - cy0 - 0.42
        rh   = max(1.10, (avail - 0.14*(len(subs)-1)) / len(subs))
        for i, q in enumerate(subs):
            ry = cy0+0.22 + i*(rh+0.14)
            sc = T["SC"][i % len(T["SC"])]
            s  = rrect(slide, rx, ry, rw, rh, T["M"], radius_pct=6)
            if s: shadow(s, blur=5, dist=2, alpha=0.10)
            rect(slide, rx, ry, rw, 0.07, sc)
            # رقم
            oval(slide, rx+0.20, ry+(rh-0.58)/2, 0.58, 0.58, sc)
            txt(slide, str(i+1), rx+0.20, ry+(rh-0.58)/2, 0.58, 0.58,
                font="Calibri", size=14, bold=True,
                color=T["D"], align=PP_ALIGN.CENTER)
            txt(slide, q, rx+0.92, ry+0.10, rw-1.06, rh-0.20,
                font=T["BF"], size=12, color=T["TL"],
                align=PP_ALIGN.RIGHT, rtl=True)
    else:
        # شريحة إشكالية فقط
        qh = H - cy0 - (2.0 if main_q else 0.40) - 0.30
        s  = rrect(slide, MX, cy0+0.20, W-MX*2, qh, T["M"], radius_pct=8)
        if s: shadow(s)
        rect(slide, MX, cy0+0.20, 0.12, qh, T["A"])
        txt(slide, "❝", MX+0.22, cy0+0.22, 1.60, 1.0,
            font="Georgia", size=52, bold=True, color=T["A"])
        txt(slide, problem, MX+0.22, cy0+1.10, W-MX*2-0.34, qh-1.22,
            font=T["BF"], size=14, color=T["TL"],
            align=PP_ALIGN.RIGHT, rtl=True)
        if main_q:
            mq_y = cy0+0.20+qh+0.18
            s2   = rrect(slide, MX, mq_y, W-MX*2, H-mq_y-0.22, T["A"], radius_pct=7)
            txt(slide, "؟", MX+0.14, mq_y+0.10, 0.80, H-mq_y-0.34,
                font="Georgia", size=48, bold=True, color=T["D"])
            txt(slide, main_q, MX+1.0, mq_y+0.14, W-MX*2-1.14, H-mq_y-0.34,
                font=T["BF"], size=15, bold=True, color=T["D"],
                align=PP_ALIGN.RIGHT, rtl=True)
    return slide


# ═══════════════════════════════════════════════════════════════════
# OBJECTIVES SLIDE — الأهداف والفرضيات
# ═══════════════════════════════════════════════════════════════════
def make_objectives(prs, data, T):
    slide  = blank(prs)
    fam    = T["FAM"]
    dark   = (fam == "NOIR")
    cy0    = header_bar(slide, T, "أهداف البحث والفرضيات",
                        "OBJECTIVES & HYPOTHESES", dark=dark)

    objs  = [o for o in data.get("objectives",[])  if o][:6]
    hypos = [h for h in data.get("hypotheses",[]) if h][:6]
    n     = max(len(objs), len(hypos), 1)
    avail = H - cy0 - 0.40
    rh    = max(1.10, (avail - 0.14*(n-1)) / n)
    cw    = (W - MX*2 - 0.32) / 2

    # ── عمود الأهداف ──────────────────────────────────
    lx = MX
    # رأس العمود
    s = rrect(slide, lx, cy0+0.18, cw, 0.56, T["A"], radius_pct=8)
    txt(slide, "🎯  الأهداف", lx, cy0+0.18, cw, 0.56,
        font=T["BF"], size=13, bold=True,
        color=T["D"], align=PP_ALIGN.CENTER)

    for i, obj in enumerate(objs):
        ry = cy0+0.86 + i*(rh+0.14)
        sc = T["SC"][i % len(T["SC"])]
        if dark:
            card_dark(slide, lx, ry, cw, rh, T, sc)
        else:
            card_light(slide, lx, ry, cw, rh, T, sc)
        # رقم دائري بارز
        nb_y = ry + (rh-0.70)/2
        oval(slide, lx+0.18, nb_y, 0.70, 0.70, sc)
        txt(slide, "%02d" % (i+1), lx+0.18, nb_y, 0.70, 0.70,
            font="Calibri", size=17, bold=True,
            color=T["D"] if dark else T["TD"], align=PP_ALIGN.CENTER)
        tc = T["TL"] if dark else T["TD"]
        txt(slide, obj, lx+1.02, ry+0.12, cw-1.16, rh-0.24,
            font=T["BF"], size=12, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)

    # ── عمود الفرضيات ─────────────────────────────────
    rx = MX + cw + 0.32
    s2 = rrect(slide, rx, cy0+0.18, cw, 0.56, T["M"], radius_pct=8)
    rect(slide, rx, cy0+0.18, cw, 0.56, T["D"] if not dark else T["M"])
    if s2:
        txt(slide, "💡  الفرضيات", rx, cy0+0.18, cw, 0.56,
            font=T["BF"], size=13, bold=True,
            color=T["A"], align=PP_ALIGN.CENTER)

    for i, hy in enumerate(hypos):
        ry = cy0+0.86 + i*(rh+0.14)
        sc = T["SC"][(i+3) % len(T["SC"])]
        s  = rrect(slide, rx, ry, cw, rh,
                   T["M"] if dark else T["CB"], radius_pct=6)
        if s: shadow(s, blur=5, dist=2, alpha=0.09)
        rect(slide, rx, ry, cw, 0.07, sc)
        # H-badge
        pill(slide, rx+cw-1.0, ry+0.10, 0.82, 0.38,
             "H%d" % (i+1), sc, T["D"] if dark else T["TD"], size=11)
        tc = T["TL"] if dark else T["TD"]
        txt(slide, hy, rx+0.18, ry+0.12, cw-1.28, rh-0.24,
            font=T["BF"], size=12, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)
    return slide


# ═══════════════════════════════════════════════════════════════════
# IMPORTANCE SLIDE — أهمية الدراسة
# ═══════════════════════════════════════════════════════════════════
def make_importance(prs, data, T):
    slide   = blank(prs)
    dark    = (T["FAM"] == "NOIR")
    cy0     = header_bar(slide, T, "أهمية الدراسة وأسباب اختيارها",
                         "RESEARCH SIGNIFICANCE", dark=dark)
    oval(slide, -3, H*0.40, W*0.55, W*0.55, T["M"])

    items = [x for x in data.get("importance",[]) if x]
    if data.get("reasons"): items.append(data["reasons"])
    items = items[:6]
    n     = len(items)
    if not items: return slide

    cols  = 2 if n > 3 else 1
    rows  = math.ceil(n/cols)
    gx, gy = 0.28, 0.22
    avail_w = W - MX*2
    avail_h = H - cy0 - 0.36
    cw = (avail_w - gx*(cols-1)) / cols
    ch = (avail_h - gy*(rows-1)) / rows
    icons = ["🔬","💡","📊","🎯","🌐","⚡"]

    for i, item in enumerate(items):
        col = i % cols
        row = i // cols
        cx  = MX + col*(cw+gx)
        cy  = cy0+0.18 + row*(ch+gy)
        sc  = T["SC"][i % len(T["SC"])]

        if dark:
            s = rrect(slide, cx, cy, cw, ch, T["M"], radius_pct=8)
        else:
            s = rrect(slide, cx, cy, cw, ch, T["CB"], radius_pct=8, line_color=T["CE"])
        if s: shadow(s, blur=7, dist=2, alpha=0.11)
        # شريط علوي + جانبي
        rect(slide, cx, cy, cw, 0.09, sc)
        rect(slide, cx, cy, 0.10, ch, sc)
        # أيقونة
        txt(slide, icons[i%len(icons)], cx+0.22, cy+0.14, 0.80, 0.72,
            font="Segoe UI Emoji", size=24, align=PP_ALIGN.CENTER)
        # رقم دائري
        oval(slide, cx+cw-0.88, cy+0.14, 0.66, 0.66, sc)
        txt(slide, "%02d"%(i+1), cx+cw-0.88, cy+0.14, 0.66, 0.66,
            font="Calibri", size=15, bold=True,
            color=T["D"] if dark else T["TD"], align=PP_ALIGN.CENTER)
        # فاصل
        line_h(slide, cx+0.22, cy+0.94, cw-0.34, sc, 0.025)
        tc = T["TL"] if dark else T["TD"]
        txt(slide, safe(item), cx+0.22, cy+1.02, cw-0.34, ch-1.14,
            font=T["BF"], size=12, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)
    return slide


# ═══════════════════════════════════════════════════════════════════
# METHODOLOGY SLIDE — المنهجية والعينة
# ═══════════════════════════════════════════════════════════════════
def make_methodology(prs, data, T):
    slide  = blank(prs)
    dark   = (T["FAM"] != "MINIMAL")
    cy0    = header_bar(slide, T, "المنهجية والعينة والمجالات",
                        "METHODOLOGY & SAMPLE", dark=dark)

    meth    = safe(data.get("methodology",""))
    stype   = safe(data.get("sampleType",""))
    ssize   = safe(data.get("sampleSize",""))
    tool_v  = safe(data.get("tool",""))
    axes    = [a for a in data.get("toolAxes",[]) if a][:4]
    spatial = safe(data.get("spatialScope",""))
    temporal= safe(data.get("temporalScope",""))
    human_s = safe(data.get("humanScope",""))
    sw      = safe(data.get("software",""))
    tests   = [t for t in data.get("statisticalTests",[]) if t][:4]

    boxes = []
    if meth:
        boxes.append(("🔬","المنهج المتبع", meth))
    if stype or ssize:
        boxes.append(("👥","العينة", " · ".join(filter(None,[stype,ssize]))))
    if tool_v:
        tv = tool_v + ("\n" + " · ".join(axes) if axes else "")
        boxes.append(("📋","أداة الدراسة", tv))
    if spatial or temporal or human_s:
        scope = "\n".join(filter(None,[
            ("📍 " + spatial) if spatial else "",
            ("🕐 " + temporal) if temporal else "",
            ("👤 " + human_s) if human_s else "",
        ]))
        boxes.append(("🌐","مجالات الدراسة", scope))
    if sw:
        swv = sw + (" · " + " · ".join(tests) if tests else "")
        boxes.append(("⚙️","البرنامج والاختبارات", swv))
    if data.get("dataSource"):
        boxes.append(("📂","مصدر البيانات", safe(data.get("dataSource",""))))

    if not boxes: return slide
    n    = len(boxes)
    cols = min(n,3)
    rows = math.ceil(n/cols)
    gx,gy = 0.26, 0.22
    bw  = (W-MX*2 - gx*(cols-1)) / cols
    bh  = (H-cy0-0.38 - gy*(rows-1)) / rows

    for i,(icon,lbl,val) in enumerate(boxes):
        col = i%cols; row = i//cols
        bx  = MX + col*(bw+gx)
        by  = cy0+0.20 + row*(bh+gy)
        sc  = T["SC"][i % len(T["SC"])]
        bg_c= T["M"] if dark else T["CB"]
        s   = rrect(slide, bx,by,bw,bh, bg_c, radius_pct=8)
        if s: shadow(s, blur=8, dist=2, alpha=0.12)
        rect(slide, bx,by,bw,0.09, sc)
        rect(slide, bx,by,0.10,bh, sc)
        # أيقونة + تسمية في صف واحد
        txt(slide, icon, bx+0.18, by+0.12, 0.72, 0.56,
            font="Segoe UI Emoji", size=20, align=PP_ALIGN.CENTER)
        tc = T["TL"] if dark else T["TD"]
        txt(slide, lbl, bx+0.96, by+0.14, bw-1.10, 0.50,
            font=T["BF"], size=13, bold=True, color=sc,
            align=PP_ALIGN.RIGHT, rtl=True)
        line_h(slide, bx+0.18, by+0.72, bw-0.28, sc, 0.025)
        txt(slide, safe(val), bx+0.18, by+0.82, bw-0.28, bh-0.94,
            font=T["BF"], size=11.5, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)
    return slide


# ═══════════════════════════════════════════════════════════════════
# KPI DASHBOARD — لوحة المؤشرات
# ═══════════════════════════════════════════════════════════════════
def make_stats(prs, data, T):
    slide  = blank(prs)
    cy0    = header_bar(slide, T,
                        "لوحة المؤشرات الإحصائية الرئيسية",
                        "KEY PERFORMANCE INDICATORS — DASHBOARD", dark=True)
    # ديكور
    oval(slide, W*0.64, cy0, W*0.46, H*0.70, T["M"])

    stats = [s for s in data.get("stats",[]) if s.get("label") and s.get("value")]
    if not stats: return slide
    n    = min(len(stats),8)
    cols = min(n,4)
    rows = math.ceil(n/cols)
    gx,gy = 0.24, 0.24
    cw  = (W-MX*2 - gx*(cols-1)) / cols
    raw_ch = (H-cy0-0.38 - gy*(rows-1)) / rows
    ch  = min(raw_ch, 3.60)
    tot = rows*ch + (rows-1)*gy
    y0  = cy0+0.20 + max(0,(H-cy0-tot-0.38)/2)

    for i,s in enumerate(stats[:8]):
        col = i%cols; row = i//cols
        cx  = MX + col*(cw+gx)
        cy  = y0  + row*(ch+gy)
        sc  = T["SC"][i % len(T["SC"])]
        kpi_card(slide, cx,cy,cw,ch, T, s["value"], s["label"], sc, dark=True)
    return slide


# ═══════════════════════════════════════════════════════════════════
# RESULTS SLIDE — النتائج
# ═══════════════════════════════════════════════════════════════════
def make_results(prs, data, T):
    slide  = blank(prs)
    dark   = (T["FAM"] != "MINIMAL")
    cy0    = header_bar(slide, T, "أهم نتائج البحث",
                        "RESEARCH FINDINGS", dark=dark)
    rect(slide, 0, 0, W, 0.10, T["A"])

    results = [r for r in data.get("mainResults",[]) if r][:7]
    if not results: return slide
    n   = len(results)
    avail = H - cy0 - 0.38
    rh  = max(0.96, (avail - 0.14*(n-1)) / n)

    for i,res in enumerate(results):
        ry  = cy0+0.20 + i*(rh+0.14)
        sc  = T["SC"][i % len(T["SC"])]
        alt = (i%2==0)
        bg_c= (T["M"] if alt else T["D"]) if dark else (T["CB"] if alt else T["L"])
        s   = rrect(slide, MX, ry, W-MX*2, rh, bg_c, radius_pct=5)
        if s: shadow(s, blur=4, dist=1.5, alpha=0.08)
        # accent bar
        rect(slide, MX, ry, W-MX*2, 0.07, sc)
        rect(slide, MX, ry, 0.10, rh, sc)
        # رقم
        oval(slide, MX+0.20, ry+(rh-0.62)/2, 0.62, 0.62, sc)
        txt(slide, "%02d"%(i+1), MX+0.20, ry+(rh-0.62)/2, 0.62, 0.62,
            font="Calibri", size=14, bold=True,
            color=T["D"] if dark else T["TD"], align=PP_ALIGN.CENTER)
        tc  = T["TL"] if dark else T["TD"]
        txt(slide, safe(res), MX+0.96, ry+0.10, W-MX*2-1.10, rh-0.20,
            font=T["BF"], size=13, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)
    return slide


# ═══════════════════════════════════════════════════════════════════
# CONCLUSION SLIDE — الخاتمة
# ═══════════════════════════════════════════════════════════════════
def make_conclusion(prs, data, T):
    slide  = blank(prs)
    cy0    = header_bar(slide, T, "الخاتمة والاستنتاجات",
                        "CONCLUSION & SYNTHESIS", dark=True)
    oval(slide, W*0.60, -2, W*0.50, H*0.75, T["M"])

    conclusion = safe(data.get("generalConclusion",""))
    recs  = [r for r in data.get("recommendations",[]) if r][:4]

    qh = H - cy0 - (2.36 if recs else 0.40) - 0.28
    s  = rrect(slide, MX, cy0+0.18, W-MX*2, qh, T["M"], radius_pct=9)
    if s: shadow(s)
    rect(slide, MX, cy0+0.18, W-MX*2, 0.09, T["A"])
    rect(slide, MX, cy0+0.18, 0.12, qh, T["A"])
    txt(slide, "❝", MX+0.22, cy0+0.20, 1.50, 1.0,
        font="Georgia", size=56, bold=True, color=T["A"])
    txt(slide, conclusion, MX+0.22, cy0+1.06, W-MX*2-0.34, qh-1.16,
        font=T["BF"], size=14, color=T["TL"],
        align=PP_ALIGN.RIGHT, rtl=True)
    txt(slide, "❞", W-1.60, cy0+qh-0.88, 1.40, 0.90,
        font="Georgia", size=56, bold=True, color=T["A"], align=PP_ALIGN.LEFT)

    # بطاقات التوصيات المختصرة
    if recs:
        ry  = H - 2.10
        line_h(slide, MX, ry, W-MX*2, T["A"], 0.05)
        rh  = 1.84
        rw  = (W-MX*2 - 0.24*(len(recs)-1)) / len(recs)
        for i,rec in enumerate(recs):
            rx = MX + i*(rw+0.24)
            sc = T["SC"][i % len(T["SC"])]
            s2 = rrect(slide, rx, ry+0.16, rw, rh, T["D"], radius_pct=7)
            if s2: shadow(s2, blur=6, dist=2, alpha=0.12)
            rect(slide, rx, ry+0.16, rw, 0.08, sc)
            txt(slide, "%02d"%(i+1), rx+0.14, ry+0.28, rw-0.28, 0.52,
                font="Calibri", size=18, bold=True,
                color=sc, align=PP_ALIGN.RIGHT)
            txt(slide, safe(rec), rx+0.14, ry+0.82, rw-0.28, rh-0.94,
                font=T["BF"], size=11, color=T["TL"],
                align=PP_ALIGN.RIGHT, rtl=True)
    return slide


# ═══════════════════════════════════════════════════════════════════
# RECOMMENDATIONS SLIDE — التوصيات
# ═══════════════════════════════════════════════════════════════════
def make_recommendations(prs, data, T):
    slide  = blank(prs)
    dark   = (T["FAM"] != "MINIMAL")
    cy0    = header_bar(slide, T, "توصيات البحث",
                        "RECOMMENDATIONS", dark=dark)

    recs = [r for r in data.get("recommendations",[]) if r][:6]
    if not recs: return slide
    n    = len(recs)
    cols = 2 if n > 3 else 1
    rows = math.ceil(n/cols)
    gx,gy = 0.30, 0.22
    cw  = (W-MX*2 - gx*(cols-1)) / cols
    ch  = (H-cy0-0.40 - gy*(rows-1)) / rows

    for i,rec in enumerate(recs):
        col = i%cols; row = i//cols
        cx  = MX + col*(cw+gx)
        cy  = cy0+0.20 + row*(ch+gy)
        sc  = T["SC"][i % len(T["SC"])]
        bg_c= T["M"] if dark else T["CB"]
        s   = rrect(slide, cx,cy,cw,ch, bg_c, radius_pct=8)
        if s: shadow(s, blur=7, dist=2, alpha=0.11)
        rect(slide, cx,cy,0.12,ch, sc)
        rect(slide, cx,cy,cw,0.09, sc)
        # رقم
        oval(slide, cx+0.22, cy+(ch-0.74)/2, 0.74, 0.74, sc)
        txt(slide, "%02d"%(i+1), cx+0.22, cy+(ch-0.74)/2, 0.74, 0.74,
            font="Calibri", size=18, bold=True,
            color=T["D"] if dark else T["TD"], align=PP_ALIGN.CENTER)
        line_v(slide, cx+1.10, cy+0.16, ch-0.32, T["CE"] if dark else T["CE"], 0.03)
        tc = T["TL"] if dark else T["TD"]
        txt(slide, safe(rec), cx+1.24, cy+0.14, cw-1.36, ch-0.28,
            font=T["BF"], size=12.5, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)
    return slide


# ═══════════════════════════════════════════════════════════════════
# FUTURE SLIDE — الآفاق
# ═══════════════════════════════════════════════════════════════════
def make_future(prs, data, T):
    slide  = blank(prs)
    cy0    = header_bar(slide, T, "آفاق البحث المستقبلية",
                        "FUTURE PERSPECTIVES", dark=True)
    items  = [f for f in data.get("futureWork",[]) if f][:5]
    if not items: return slide
    n      = len(items)
    avail  = H - cy0 - 0.40
    rh     = max(1.0, (avail - 0.16*(n-1)) / n)

    for i,fut in enumerate(items):
        ry = cy0+0.22 + i*(rh+0.16)
        sc = T["SC"][i % len(T["SC"])]
        s  = rrect(slide, MX,ry, W-MX*2, rh, T["M"], radius_pct=6)
        if s: shadow(s, blur=5, dist=1.5, alpha=0.09)
        rect(slide, MX,ry, 0.12,rh, sc)
        txt(slide, "🚀", MX+0.22, ry+(rh-0.55)/2, 0.55, 0.55,
            font="Segoe UI Emoji", size=20, align=PP_ALIGN.CENTER)
        oval(slide, W-MX-0.76, ry+(rh-0.56)/2, 0.56, 0.56, sc)
        txt(slide, str(i+1), W-MX-0.76, ry+(rh-0.56)/2, 0.56, 0.56,
            font="Calibri", size=14, bold=True,
            color=T["D"], align=PP_ALIGN.CENTER)
        txt(slide, safe(fut), MX+0.88, ry+0.12, W-MX*2-1.72, rh-0.24,
            font=T["BF"], size=13, color=T["TL"],
            align=PP_ALIGN.RIGHT, rtl=True)
    return slide


# ═══════════════════════════════════════════════════════════════════
# REFERENCES SLIDE — المراجع
# ═══════════════════════════════════════════════════════════════════
def make_references(prs, data, T):
    refs = [r for r in data.get("references",[]) if r][:6]
    if not refs: return
    slide  = blank(prs)
    dark   = (T["FAM"] == "NOIR")
    cy0    = header_bar(slide, T, "أبرز المراجع والمصادر",
                        "KEY REFERENCES", dark=dark)
    n   = len(refs)
    avail = H - cy0 - 0.38
    rh  = max(0.92, (avail - 0.12*(n-1)) / n)

    for i,ref in enumerate(refs):
        ry  = cy0+0.20 + i*(rh+0.12)
        sc  = T["SC"][i % len(T["SC"])]
        bg_c= T["M"] if (dark and i%2==0) else (T["D"] if dark else (T["CB"] if i%2==0 else T["L"]))
        s   = rrect(slide, MX,ry, W-MX*2, rh, bg_c, radius_pct=5)
        rect(slide, MX,ry, 0.10,rh, sc)
        oval(slide, MX+0.18, ry+(rh-0.56)/2, 0.56, 0.56, sc)
        txt(slide, str(i+1), MX+0.18, ry+(rh-0.56)/2, 0.56, 0.56,
            font="Calibri", size=14, bold=True,
            color=T["D"] if dark else T["TD"], align=PP_ALIGN.CENTER)
        tc = T["TL"] if dark else T["TD"]
        txt(slide, safe(ref), MX+0.88, ry+0.10, W-MX*2-1.00, rh-0.20,
            font=T["BF"], size=11, color=tc,
            align=PP_ALIGN.RIGHT, rtl=True)
    return slide


# ═══════════════════════════════════════════════════════════════════
# THANK YOU SLIDE — شريحة الشكر
# ═══════════════════════════════════════════════════════════════════
def make_final(prs, data, T):
    slide = blank(prs)
    bg(slide, T["D"])
    # خلفية هندسية ضخمة
    oval(slide, W*0.55, -4.0, W*0.62, H*0.95, T["M"])
    oval(slide, -5.0, H*0.48, W*0.72, W*0.72, T["A"])
    oval(slide, W*0.30, H*0.20, W*0.40, W*0.40, T["M"])

    rect(slide, 0, 0, W, 0.65, T["A"])
    rect(slide, 0, H-0.65, W, 0.65, T["A"])

    # نص الشكر الرئيسي (ثلاث لغات)
    txt(slide, "شكراً لحسن استماعكم", MX, H*0.20, W-MX*2, 1.70,
        font=T["HF"], size=44, bold=True,
        color=T["TL"], align=PP_ALIGN.CENTER)
    txt(slide, "Merci pour votre attention", MX, H*0.20+1.80, W-MX*2, 0.82,
        font="Calibri", size=22, italic=True,
        color=T["A"], align=PP_ALIGN.CENTER, rtl=False)
    txt(slide, "Thank you for your kind attention", MX, H*0.20+2.68, W-MX*2, 0.56,
        font="Calibri", size=14, italic=True,
        color=T["TM"], align=PP_ALIGN.CENTER, rtl=False)

    # فاصل مزخرف
    line_h(slide, W*0.30, H*0.68, W*0.40, T["A"], 0.06)
    oval(slide, W/2-0.20, H*0.68-0.14, 0.40, 0.40, T["A"])

    # بيانات الطالب
    student  = safe(data.get("studentName",""))
    sup      = safe(data.get("supervisor",""))
    yr       = safe(data.get("year",""))
    info     = "  ·  ".join(filter(None,["إعداد: "+student if student else "", "إشراف: "+sup if sup else ""]))
    if info:
        txt(slide, info, MX, H*0.73, W-MX*2, 0.56,
            font=T["BF"], size=13, color=T["TM"],
            align=PP_ALIGN.CENTER, rtl=True)
    univ = safe(data.get("university",""))
    if univ:
        txt(slide, univ, MX, H*0.80, W-MX*2, 0.48,
            font=T["BF"], size=11, italic=True,
            color=T["A"], align=PP_ALIGN.CENTER, rtl=True)
    if yr:
        pill(slide, W/2-1.10, H-0.48, 2.20, 0.38,
             yr, T["A"], T["D"], size=11)
    return slide


# ═══════════════════════════════════════════════════════════════════
# ORCHESTRATOR
# ═══════════════════════════════════════════════════════════════════
def generate_presentation(data: dict, output_path: str) -> None:
    key = data.get("theme","navy_gold")
    T   = PALETTES.get(key, PALETTES["navy_gold"])

    prs = Presentation()
    prs.slide_width  = Cm(W)
    prs.slide_height = Cm(H)

    cfg   = data.get("slides",{})
    def show(k): return cfg.get(k,True)
    def fl(k):   return [x for x in data.get(k,[]) if x]

    make_cover(prs, data, T)

    if show("intro") and (data.get("introOverview") or data.get("introApproach")):
        make_intro(prs, data, T)

    chs = [c for c in data.get("chapters",[]) if c.get("title")]
    if show("plan") and chs:
        make_plan(prs, data, T, chs)

    if show("problem") and (data.get("mainProblem") or data.get("mainQuestion") or fl("subQuestions")):
        make_problem(prs, data, T)

    if show("objectives") and (fl("objectives") or fl("hypotheses")):
        make_objectives(prs, data, T)

    if show("importance") and (fl("importance") or data.get("reasons")):
        make_importance(prs, data, T)

    if show("methodology") and (data.get("methodology") or data.get("sampleType") or data.get("tool")):
        make_methodology(prs, data, T)

    stats = [s for s in data.get("stats",[]) if s.get("label") and s.get("value")]
    if show("kpi") and stats:
        make_stats(prs, data, T)

    if show("results") and fl("mainResults"):
        make_results(prs, data, T)

    if show("conclusion") and data.get("generalConclusion"):
        make_conclusion(prs, data, T)

    if show("recommendations") and fl("recommendations"):
        make_recommendations(prs, data, T)

    if show("future") and fl("futureWork"):
        make_future(prs, data, T)

    if show("references") and fl("references"):
        make_references(prs, data, T)

    if show("thankyou"):
        make_final(prs, data, T)

    prs.save(output_path)
    n = len(prs.slides._sldIdLst)
    print("✅  %d slides [canva·%s·%s] → %s" % (n, T["FAM"], key, output_path),
          file=sys.stderr)


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python generator_canva.py input.json output.pptx", file=sys.stderr)
        sys.exit(1)
    with open(sys.argv[1], encoding="utf-8") as f:
        payload = json.load(f)
    generate_presentation(payload, sys.argv[2])
