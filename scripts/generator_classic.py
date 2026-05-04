import math
"""
مذكرتي Pro — محرك التحويل البصري v5.0
Radical Design System · 3 Structural Layouts × 8 Palettes
Every layout is architecturally different — not just colours.
"""
import json, sys, datetime, math
from pptx import Presentation
from pptx.util import Pt, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W, H = 33.867, 19.05   # cm – Widescreen 16:9

# ─────────────────────────────────────────────────────────────────────
# HELPERS
# ─────────────────────────────────────────────────────────────────────
def cm(v): return Cm(v)
def rgb(r, g, b): return RGBColor(r, g, b)
def safe(v, fb=""): return str(v).strip() if v else fb
def clamp(v, lo, hi): return max(lo, min(hi, v))

def blank(prs):   return prs.slides.add_slide(prs.slide_layouts[6])
def bg(s, c):
    f = s.background.fill; f.solid(); f.fore_color.rgb = c

def rect(slide, x, y, w, h, fill, line_color=None, line_pt=0.5, alpha=0):
    if w <= 0 or h <= 0: return None
    sp = slide.shapes.add_shape(1, cm(x), cm(y), cm(w), cm(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if alpha: sp.fill.fore_color.transparency = clamp(alpha, 0.0, 0.99)
    if line_color:
        sp.line.color.rgb = line_color; sp.line.width = Pt(line_pt)
    else: sp.line.fill.background()
    return sp

def oval(slide, x, y, w, h, fill, alpha=0):
    if w <= 0 or h <= 0: return None
    sp = slide.shapes.add_shape(9, cm(x), cm(y), cm(w), cm(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if alpha: sp.fill.fore_color.transparency = clamp(alpha, 0.0, 0.99)
    sp.line.fill.background(); return sp

def txt(slide, text, x, y, w, h,
        font="Cairo", size=13, bold=False, italic=False,
        color=None, align=PP_ALIGN.RIGHT, mg=0.1):
    if w <= 0 or h <= 0: return None
    tb = slide.shapes.add_textbox(cm(x), cm(y), cm(w), cm(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = cm(mg); tf.margin_right = cm(mg)
    tf.margin_top = cm(0.04); tf.margin_bottom = cm(0.04)
    p = tf.paragraphs[0]; p.alignment = align
    run = p.add_run()
    run.text = str(text) if text is not None else ""
    run.font.name = font; run.font.size = Pt(size)
    run.font.bold = bold; run.font.italic = italic
    if color: run.font.color.rgb = color
    return tb

def ln(slide, x, y, w, color, h=0.06):
    rect(slide, x, y, w, h, color)

# ─────────────────────────────────────────────────────────────────────
# TYPOGRAPHY SYSTEM
# ─────────────────────────────────────────────────────────────────────
class TY:
    DISPLAY   = dict(size=26, bold=True,  italic=False)
    H1        = dict(size=28, bold=True,  italic=False)
    H2        = dict(size=16, bold=True,  italic=False)
    H3        = dict(size=13, bold=True,  italic=False)
    OVERLINE  = dict(size=9,  bold=True,  italic=False)
    BODY      = dict(size=13, bold=False, italic=False)
    BODY_SM   = dict(size=11, bold=False, italic=False)
    BULLET    = dict(size=12, bold=False, italic=False)
    KPI_VAL_MAX = 56
    KPI_VAL_MIN = 28
    KPI_LBL   = dict(size=12, bold=False, italic=False)
    CAPTION   = dict(size=9,  bold=False, italic=True)
    LABEL     = dict(size=13, bold=True,  italic=False)
    META      = dict(size=11, bold=False, italic=False)
    EN_SUB    = dict(size=11, bold=False, italic=True)
    NUMBER    = dict(size=28, bold=True,  italic=False)



# ─────────────────────────────────────────────────────────────────────
# PALETTES  (8 colour sets — architecture-agnostic)
# Keys: D=dark  M=mid  L=light  A=accent  A2=accent-light
#       TL=text-light  TD=text-dark  TM=text-muted
#       CB=card-bg  CE=card-edge  HF=heading-font  BF=body-font
#       SC=stripe-colours[6]  LAYOUT=which design system to use
# ─────────────────────────────────────────────────────────────────────
PALETTES = {
    "navy_gold": dict(
        D=rgb(0x07,0x17,0x2F), M=rgb(0x0E,0x27,0x4D), L=rgb(0xF4,0xF6,0xFB),
        A=rgb(0xC6,0xA0,0x3C), A2=rgb(0xE8,0xC9,0x7B),
        TL=rgb(0xFF,0xFF,0xFF), TD=rgb(0x07,0x17,0x2F), TM=rgb(0x64,0x74,0x8B),
        CB=rgb(0xFF,0xFF,0xFF), CE=rgb(0xE2,0xE8,0xF0),
        HF="Georgia", BF="Cairo",
        SC=[rgb(0xC6,0xA0,0x3C),rgb(0x1A,0x40,0x72),rgb(0xE8,0xC9,0x7B),
            rgb(0x2A,0x55,0x98),rgb(0x0E,0x27,0x4D),rgb(0x8A,0x6E,0x28)],
        LAYOUT="classic",
    ),
    "dark_teal": dict(
        D=rgb(0x06,0x1A,0x28), M=rgb(0x04,0x4E,0x6E), L=rgb(0xEF,0xFD,0xFA),
        A=rgb(0x00,0xD4,0xAA), A2=rgb(0x67,0xE8,0xD3),
        TL=rgb(0xFF,0xFF,0xFF), TD=rgb(0x06,0x1A,0x28), TM=rgb(0x52,0x73,0x84),
        CB=rgb(0xFF,0xFF,0xFF), CE=rgb(0xCC,0xF5,0xED),
        HF="Trebuchet MS", BF="Cairo",
        SC=[rgb(0x00,0xD4,0xAA),rgb(0x04,0x4E,0x6E),rgb(0x09,0x93,0xC3),
            rgb(0x67,0xE8,0xD3),rgb(0x03,0x30,0x44),rgb(0x06,0x7A,0x9F)],
        LAYOUT="bold",
    ),
    "burgundy": dict(
        D=rgb(0x3A,0x00,0x18), M=rgb(0x6B,0x15,0x37), L=rgb(0xFD,0xF0,0xF4),
        A=rgb(0xF0,0xB8,0xCC), A2=rgb(0xFA,0xD4,0xE2),
        TL=rgb(0xFF,0xFF,0xFF), TD=rgb(0x3A,0x00,0x18), TM=rgb(0x78,0x55,0x63),
        CB=rgb(0xFF,0xFF,0xFF), CE=rgb(0xF5,0xD8,0xE4),
        HF="Georgia", BF="Cairo",
        SC=[rgb(0xF0,0xB8,0xCC),rgb(0x6B,0x15,0x37),rgb(0xCC,0x3D,0x73),
            rgb(0xFA,0xD4,0xE2),rgb(0x3A,0x00,0x18),rgb(0x9A,0x20,0x50)],
        LAYOUT="minimal",
    ),
    "forest": dict(
        D=rgb(0x0F,0x2D,0x1E), M=rgb(0x1E,0x4D,0x36), L=rgb(0xF0,0xFB,0xF4),
        A=rgb(0x86,0xBB,0x56), A2=rgb(0xB4,0xD9,0x84),
        TL=rgb(0xFF,0xFF,0xFF), TD=rgb(0x0F,0x2D,0x1E), TM=rgb(0x4A,0x6B,0x56),
        CB=rgb(0xFF,0xFF,0xFF), CE=rgb(0xD1,0xF0,0xDA),
        HF="Cambria", BF="Cairo",
        SC=[rgb(0x86,0xBB,0x56),rgb(0x1E,0x4D,0x36),rgb(0x4A,0x99,0x6B),
            rgb(0xB4,0xD9,0x84),rgb(0x0F,0x2D,0x1E),rgb(0x2E,0x7A,0x56)],
        LAYOUT="classic",
    ),
    "midnight_purple": dict(
        D=rgb(0x12,0x05,0x2E), M=rgb(0x2D,0x10,0x6B), L=rgb(0xF5,0xF0,0xFF),
        A=rgb(0xC0,0x7A,0xFF), A2=rgb(0xDD,0xB3,0xFF),
        TL=rgb(0xFF,0xFF,0xFF), TD=rgb(0x12,0x05,0x2E), TM=rgb(0x7A,0x6A,0x95),
        CB=rgb(0xFF,0xFF,0xFF), CE=rgb(0xE8,0xD8,0xFF),
        HF="Georgia", BF="Cairo",
        SC=[rgb(0xC0,0x7A,0xFF),rgb(0x2D,0x10,0x6B),rgb(0x7B,0x3F,0xE0),
            rgb(0xDD,0xB3,0xFF),rgb(0x12,0x05,0x2E),rgb(0x52,0x1A,0xC1)],
        LAYOUT="bold",
    ),
    "charcoal_orange": dict(
        D=rgb(0x1A,0x1A,0x2E), M=rgb(0x2D,0x2D,0x44), L=rgb(0xFF,0xF8,0xF2),
        A=rgb(0xFF,0x6B,0x35), A2=rgb(0xFF,0xA0,0x70),
        TL=rgb(0xFF,0xFF,0xFF), TD=rgb(0x1A,0x1A,0x2E), TM=rgb(0x7A,0x7A,0x95),
        CB=rgb(0xFF,0xFF,0xFF), CE=rgb(0xFF,0xE5,0xD5),
        HF="Trebuchet MS", BF="Cairo",
        SC=[rgb(0xFF,0x6B,0x35),rgb(0x2D,0x2D,0x44),rgb(0xFF,0xA0,0x70),
            rgb(0xE5,0x4E,0x1A),rgb(0x1A,0x1A,0x2E),rgb(0xFF,0xC2,0x9E)],
        LAYOUT="minimal",
    ),
    "ice_blue": dict(
        D=rgb(0x0A,0x25,0x4A), M=rgb(0x1A,0x4A,0x8A), L=rgb(0xF0,0xF6,0xFF),
        A=rgb(0x4A,0xB3,0xFF), A2=rgb(0x8E,0xD0,0xFF),
        TL=rgb(0xFF,0xFF,0xFF), TD=rgb(0x0A,0x25,0x4A), TM=rgb(0x55,0x77,0xAA),
        CB=rgb(0xFF,0xFF,0xFF), CE=rgb(0xD0,0xE8,0xFF),
        HF="Calibri", BF="Cairo",
        SC=[rgb(0x4A,0xB3,0xFF),rgb(0x1A,0x4A,0x8A),rgb(0x8E,0xD0,0xFF),
            rgb(0x0A,0x25,0x4A),rgb(0x23,0x72,0xD9),rgb(0xBD,0xE3,0xFF)],
        LAYOUT="classic",
    ),
    "sand_gold": dict(
        D=rgb(0x2C,0x1A,0x0E), M=rgb(0x5C,0x3D,0x20), L=rgb(0xFD,0xF8,0xF0),
        A=rgb(0xC8,0x86,0x1F), A2=rgb(0xE8,0xB9,0x64),
        TL=rgb(0xFF,0xFF,0xFF), TD=rgb(0x2C,0x1A,0x0E), TM=rgb(0x7A,0x62,0x4A),
        CB=rgb(0xFF,0xFF,0xFF), CE=rgb(0xF0,0xE0,0xC8),
        HF="Georgia", BF="Cairo",
        SC=[rgb(0xC8,0x86,0x1F),rgb(0x5C,0x3D,0x20),rgb(0xE8,0xB9,0x64),
            rgb(0x2C,0x1A,0x0E),rgb(0x8A,0x5E,0x30),rgb(0xF0,0xD0,0x88)],
        LAYOUT="bold",
    ),
}

# ═════════════════════════════════════════════════════════════════════
# ██████████████  LAYOUT  CLASSIC  ████████████████████████████████████
# Architecture: dark full-bleed slides alternating with light slides
# Visual DNA: full-width accent bars · left-side pillar · shadow cards
# ═════════════════════════════════════════════════════════════════════

class Classic:
    """Standard dark/light alternating layout — refined & authoritative"""

    @staticmethod
    def cover(slide, T, data):
        bg(slide, T["D"])
        # ── ديكور هندسي خلفي ──────────────────────────────
        oval(slide, W*0.55, -3,   W*0.55, H*0.80, T["M"], 0.32)
        oval(slide, W*0.72, -1,   W*0.35, H*0.60, T["A"], 0.72)
        oval(slide, -4,    H*0.65, 11,    11,      T["M"], 0.52)
        # ── شريط accent رأسي + سفلي ───────────────────────
        rect(slide, W-0.65, 0, 0.65, H, T["A"])
        rect(slide, 0, H-0.50, W-0.65, 0.50, T["A"])
        # ── نطاق الجامعة ──────────────────────────────────
        rect(slide, 0, 0, W-0.65, 2.60, T["M"])
        # ── خط فاصل ذهبي ──────────────────────────────────
        rect(slide, 0, 2.60, W-0.65, 0.055, T["A"])
        # ── رقم ديكوري كبير خفي ───────────────────────────
        yr = safe(data.get("year","")).replace("–","-").split("-")
        yr_txt = yr[-1] if yr else ""
        txt(slide, yr_txt, W*0.02, H*0.28, W*0.48, H*0.55,
            font="Calibri", size=200, bold=True,
            color=T["M"], align=PP_ALIGN.LEFT)
        # ── شريط كلمات مفتاحية أسفل ──────────────────────
        kw = safe(data.get("keywords",""))
        if kw:
            rect(slide, 0, H-0.50, W-0.65, 0.50, T["A"])
            txt(slide, kw, 0.22, H-0.48, W-1.10, 0.40,
                font=T["BF"], size=9, italic=True,
                color=T["D"], align=PP_ALIGN.RIGHT)
        _cover_shared(slide, T, data)

    @staticmethod
    def toc(slide, T, chapters):
        bg(slide, T["L"])
        rect(slide, 0, 0, 0.52, H, T["D"])
        rect(slide, 0.52, 0, 0.20, H, T["A"])
        rect(slide, 0.72, 0, W-0.72, 3.08, T["D"])
        txt(slide, "المحتويات", W-1.72, 0.40, W-2.45, 1.48,
            font=T["HF"], size=34, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
        txt(slide, "Table of Contents", 1.38, 1.94, W-2.75, 0.68,
            font="Calibri", size=14, italic=True, color=T["A"], align=PP_ALIGN.LEFT)
        cw = (W-3.08)/2; gx = 0.38
        ch = (H-4.0)/3 - 0.20
        for i, ch_item in enumerate(chapters[:6]):
            col = i%2; row = i//2
            cx = 0.90 + col*(cw+gx); cy = 3.42 + row*(ch+0.22)
            sc = T["SC"][i%len(T["SC"])]
            rect(slide, cx+0.10, cy+0.10, cw, ch, T["CE"])
            rect(slide, cx, cy, cw, ch, T["CB"], line_color=T["CE"])
            rect(slide, cx, cy, 0.19, ch, sc)
            txt(slide, f"{i+1:02d}", cx+0.28, cy+0.10, 1.58, ch*0.50,
                font="Calibri", size=30, bold=True, color=sc, align=PP_ALIGN.LEFT)
            txt(slide, ch_item["title"], cx+0.28, cy+ch*0.52, cw-0.46, ch*0.36,
                font=T["BF"], size=14, bold=True, color=T["TD"], align=PP_ALIGN.RIGHT)
            if ch_item.get("sub"):
                txt(slide, ch_item["sub"], cx+0.28, cy+ch*0.84, cw-0.46, 0.44,
                    font="Calibri", size=10, italic=True, color=T["TM"], align=PP_ALIGN.LEFT)

    @staticmethod
    def section_dark(slide, T, title_ar, sub_en=None):
        """Full-width dark slide header. Returns content_y."""
        rect(slide, 0, 0, W, 0.40, T["A"])
        txt(slide, title_ar, 1.0, 0.60, W-2.0, 1.16,
            font=T["HF"], size=28, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
        if sub_en:
            txt(slide, sub_en, 1.0, 1.82, W-2.0, 0.60,
                font="Calibri", size=12, italic=True, color=T["A"], align=PP_ALIGN.RIGHT)
        return 2.64

    @staticmethod
    def section_light(slide, T, title_ar, sub_en=None):
        """Light slide header with left pillar. Returns content_y."""
        rect(slide, 0, 0, W, 0.40, T["D"])
        txt(slide, title_ar, 1.0, 0.60, W-2.0, 1.16,
            font=T["HF"], size=28, bold=True, color=T["TD"], align=PP_ALIGN.RIGHT)
        if sub_en:
            txt(slide, sub_en, 1.0, 1.82, W-2.0, 0.60,
                font="Calibri", size=12, italic=True, color=T["TM"], align=PP_ALIGN.RIGHT)
        ln(slide, 1.0, 2.62, W-2.0, T["A"])
        return 2.78

    @staticmethod
    def kpi_card(slide, T, x, y, w, h, value, label, sc):
        ac = sc or T["A"]
        rect(slide, x, y, w, h, T["M"])
        rect(slide, x, y,     w, 0.28, ac)
        rect(slide, x, y+h-0.28, w, 0.28, ac)
        _kpi_text(slide, T, x, y, w, h, value, label, ac)

    @staticmethod
    def result_row(slide, T, x, y, w, h, i, text):
        sc = T["SC"][i%len(T["SC"])]
        rect(slide, x+0.08, y+0.06, w, h, T["CE"])
        rect(slide, x,     y,    w, h, T["CB"], line_color=T["CE"])
        rect(slide, x,     y,    0.24, h, sc)
        nw = 1.65; rect(slide, x, y, nw, h, T["D"])
        txt(slide, str(i+1), x, y+h/2-0.65, nw, 1.3,
            font="Calibri", size=28, bold=True, color=sc, align=PP_ALIGN.CENTER)
        bs = min(h-0.18, 0.92); bx = x+w-bs-0.18; by = y+(h-bs)/2
        rect(slide, bx, by, bs, bs, T["D"])
        txt(slide, "✓", bx, by, bs, bs,
            font="Calibri", size=14, bold=True, color=sc, align=PP_ALIGN.CENTER)
        txt(slide, text, x+nw+0.16, y+0.10, bx-x-nw-0.28, h-0.20,
            font=T["BF"], size=12, color=T["TD"], align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════
# ██████████████  LAYOUT  BOLD  ███████████████████████████████████████
# Architecture: high-contrast blocks · thick accent panels · oversized type
# Visual DNA: left accent column · full-bleed colour zones · big numbers
# ═════════════════════════════════════════════════════════════════════

class Bold:
    """High-impact layout — large typography, thick bands, strong contrast"""

    @staticmethod
    def cover(slide, T, data):
        bg(slide, T["D"])
        # Bold left accent pillar
        rect(slide, 0, 0, 1.8, H, T["A"])
        # Large decorative circle
        oval(slide, W*0.50, -4, H*1.1, H*1.1, T["M"], 0.42)
        oval(slide, W*0.68,  2, H*0.65, H*0.65, T["A"], 0.82)
        # Top-right corner accent
        rect(slide, W-4.5, 0, 4.5, 0.7, T["A"])
        # University band
        rect(slide, 1.8, 0, W-1.8, 2.52, T["M"])
        _cover_shared(slide, T, data, offset_x=1.9, max_x=W*0.72)

    @staticmethod
    def toc(slide, T, chapters):
        bg(slide, T["D"])
        rect(slide, 0, 0, 1.8, H, T["A"])
        rect(slide, 1.8, 0, W-1.8, H, T["D"])
        txt(slide, "المحتويات", W-1.5, 0.5, W-3.2, 1.8,
            font=T["HF"], size=40, bold=True, color=T["A"], align=PP_ALIGN.RIGHT)
        txt(slide, "Table of Contents", 2.2, 2.3, 10, 0.75,
            font="Calibri", size=15, italic=True, color=T["TM"], align=PP_ALIGN.LEFT)
        ln(slide, 2.0, 3.25, W-3.0, T["A"], 0.08)
        # Horizontal 3-col grid — fixed height cards
        cw = (W-3.8)/3; ch = 3.5; gapx = 0.4; gapy = 0.4
        for i, ch_item in enumerate(chapters[:6]):
            col = i%3; row = i//3
            cx = 2.0 + col*(cw+gapx)
            cy = 3.6 + row*(ch+gapy)
            sc = T["SC"][i%len(T["SC"])]
            rect(slide, cx, cy, cw, ch, T["M"])
            rect(slide, cx, cy, cw, 0.42, sc)
            # Number top-left
            txt(slide, f"{i+1:02d}", cx+0.18, cy+0.06, 1.5, 0.36,
                font="Calibri", size=16, bold=True, color=T["D"], align=PP_ALIGN.LEFT)
            # Arabic title — centered in card
            txt(slide, ch_item["title"], cx+0.16, cy+0.62, cw-0.32, 1.6,
                font=T["BF"], size=14, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
            # EN sub — bottom
            if ch_item.get("sub"):
                txt(slide, ch_item["sub"], cx+0.16, cy+ch-0.72, cw-0.32, 0.56,
                    font="Calibri", size=10, italic=True, color=T["A"], align=PP_ALIGN.LEFT)

    @staticmethod
    def section_dark(slide, T, title_ar, sub_en=None):
        rect(slide, 0, 0, 1.8, H, T["A"])
        txt(slide, title_ar, 2.2, 0.55, W-3.2, 1.2,
            font=T["HF"], size=28, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
        if sub_en:
            txt(slide, sub_en, 2.2, 1.82, W-3.2, 0.60,
                font="Calibri", size=12, italic=True, color=T["A"], align=PP_ALIGN.RIGHT)
        ln(slide, 2.2, 2.60, W-3.2, T["A"])
        return 2.76

    @staticmethod
    def section_light(slide, T, title_ar, sub_en=None):
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
        txt(slide, title_ar, 2.2, 0.55, W-3.2, 1.2,
            font=T["HF"], size=28, bold=True, color=T["TD"], align=PP_ALIGN.RIGHT)
        if sub_en:
            txt(slide, sub_en, 2.2, 1.82, W-3.2, 0.60,
                font="Calibri", size=12, italic=True, color=T["TM"], align=PP_ALIGN.RIGHT)
        ln(slide, 2.2, 2.60, W-3.2, T["A"])
        return 2.76

    @staticmethod
    def kpi_card(slide, T, x, y, w, h, value, label, sc):
        ac = sc or T["A"]
        rect(slide, x, y, w, h, T["M"])
        rect(slide, x, y, 0.22, h, ac)          # left accent bar
        rect(slide, x, y+h-0.28, w, 0.28, ac)   # bottom bar
        _kpi_text(slide, T, x+0.22, y, w-0.22, h, value, label, ac)

    @staticmethod
    def result_row(slide, T, x, y, w, h, i, text):
        sc = T["SC"][i%len(T["SC"])]
        rect(slide, x, y, w, h, T["M"] if i%2==0 else T["D"])
        rect(slide, x, y, 0.30, h, sc)
        idx_w = 1.1
        txt(slide, f"{i+1:02d}", x+0.38, y+h/2-0.52, idx_w, 1.05,
            font="Calibri", size=22, bold=True, color=sc, align=PP_ALIGN.RIGHT)
        txt(slide, text, x+idx_w+0.55, y+0.10, w-idx_w-0.70, h-0.20,
            font=T["BF"], size=12, color=T["TL"], align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════
# ██████████████  LAYOUT  MINIMAL  ████████████████████████████████████
# Architecture: white/light space-forward · thin lines · elegant type
# Visual DNA: bottom-anchored accents · generous whitespace · subtle tones
# ═════════════════════════════════════════════════════════════════════

class Minimal:
    """Editorial minimal layout — whitespace-forward, thin lines, refined"""

    @staticmethod
    def cover(slide, T, data):
        bg(slide, T["L"])
        # Large decorative circle — pushed far right, no text overlap
        oval(slide, W*0.60, -2, H*1.05, H*1.05, T["CE"], 0.0)
        oval(slide, W*0.75, H*0.30, H*0.65, H*0.65, T["A"], 0.85)
        # Left dark content zone
        rect(slide, 0, 0, W*0.60, H, T["D"])
        # Thin top + bottom accent lines
        rect(slide, 0, 0, W, 0.30, T["A"])
        rect(slide, 0, H-0.42, W, 0.42, T["D"])
        # University band stays on dark zone
        rect(slide, 0, 0, W*0.60, 2.52, T["M"])
        _cover_shared(slide, T, data, offset_x=0.38, max_x=W*0.60-0.3)

    @staticmethod
    def toc(slide, T, chapters):
        bg(slide, T["L"])
        rect(slide, 0, 0, W, 0.28, T["D"])
        rect(slide, 0, H-0.28, W, 0.28, T["A"])
        txt(slide, "المحتويات", W-1.5, 0.55, W-2.5, 1.6,
            font=T["HF"], size=36, bold=True, color=T["TD"], align=PP_ALIGN.RIGHT)
        txt(slide, "Table of Contents", 1.4, 1.95, 12, 0.7,
            font="Calibri", size=14, italic=True, color=T["TM"], align=PP_ALIGN.LEFT)
        ln(slide, 1.4, 2.88, W-2.8, T["A"], 0.05)
        # Clean numbered list
        cw = (W-2.8)/2 - 0.3; gy = (H-3.6)/3 - 0.28
        for i, ch_item in enumerate(chapters[:6]):
            col = i%2; row = i//2
            cx = 1.4 + col*(cw+0.6); cy = 3.1 + row*(gy+0.28)
            sc = T["SC"][i%len(T["SC"])]
            ln(slide, cx, cy, cw, sc, 0.12)
            txt(slide, f"{i+1:02d}", cx, cy+0.20, 1.4, gy*0.52,
                font="Calibri", size=28, bold=True, color=sc, align=PP_ALIGN.LEFT)
            txt(slide, ch_item["title"], cx+1.5, cy+0.22, cw-1.5, gy*0.48,
                font=T["BF"], size=14, bold=True, color=T["TD"], align=PP_ALIGN.RIGHT)
            if ch_item.get("sub"):
                txt(slide, ch_item["sub"], cx+1.5, cy+gy*0.58, cw-1.5, 0.48,
                    font="Calibri", size=10, italic=True, color=T["TM"], align=PP_ALIGN.RIGHT)
            ln(slide, cx, cy+gy+0.06, cw, T["CE"], 0.05)

    @staticmethod
    def section_dark(slide, T, title_ar, sub_en=None):
        bg(slide, T["D"])
        rect(slide, 0, 0, W, 0.28, T["A"])
        txt(slide, title_ar, 1.4, 0.50, W-2.8, 1.22,
            font=T["HF"], size=28, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
        if sub_en:
            txt(slide, sub_en, 1.4, 1.80, W-2.8, 0.60,
                font="Calibri", size=12, italic=True, color=T["A"], align=PP_ALIGN.RIGHT)
        ln(slide, 1.4, 2.58, W-2.8, T["A"], 0.05)
        return 2.72

    @staticmethod
    def section_light(slide, T, title_ar, sub_en=None):
        bg(slide, T["L"])
        rect(slide, 0, 0, W, 0.28, T["A"])
        rect(slide, 0, H-0.28, W, 0.28, T["D"])
        txt(slide, title_ar, 1.4, 0.50, W-2.8, 1.22,
            font=T["HF"], size=28, bold=True, color=T["TD"], align=PP_ALIGN.RIGHT)
        if sub_en:
            txt(slide, sub_en, 1.4, 1.80, W-2.8, 0.60,
                font="Calibri", size=12, italic=True, color=T["TM"], align=PP_ALIGN.RIGHT)
        ln(slide, 1.4, 2.58, W-2.8, T["A"], 0.05)
        return 2.72

    @staticmethod
    def kpi_card(slide, T, x, y, w, h, value, label, sc):
        ac = sc or T["A"]
        rect(slide, x, y, w, h, T["CB"], line_color=T["CE"])
        ln(slide, x, y, w, ac, 0.22)             # top accent line
        ln(slide, x, y+h-0.22, w, ac, 0.22)      # bottom accent line
        _kpi_text(slide, T, x, y, w, h, value, label, ac, text_color_override=T["TD"])

    @staticmethod
    def result_row(slide, T, x, y, w, h, i, text):
        sc = T["SC"][i%len(T["SC"])]
        rect(slide, x, y, w, h, T["CB"], line_color=T["CE"])
        ln(slide, x, y, w, sc, 0.12)
        marker = f"0{i+1}" if i < 9 else str(i+1)
        txt(slide, marker, x+0.2, y+0.1, 1.2, h-0.2,
            font="Calibri", size=20, bold=True, color=sc, align=PP_ALIGN.RIGHT)
        ln(slide, x+1.5, y+h/2-0.3, 0.05, sc, 0.6)
        txt(slide, text, x+1.7, y+0.10, w-2.0, h-0.20,
            font=T["BF"], size=12, color=T["TD"], align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────
# LAYOUT REGISTRY
# ─────────────────────────────────────────────────────────────────────
LAYOUTS = {"classic": Classic, "bold": Bold, "minimal": Minimal}

def get_layout(T):
    return LAYOUTS.get(T.get("LAYOUT", "classic"), Classic)


# ─────────────────────────────────────────────────────────────────────
# SHARED KPI TEXT (reused by all three layout families)
# ─────────────────────────────────────────────────────────────────────
def _kpi_text(slide, T, x, y, w, h, value, label, ac, text_color_override=None):
    tc = text_color_override or T["TL"]
    interior = h - 0.58
    vs  = clamp(58 - max(0, len(str(value))-4)*5, 28, 56)
    lhv = vs * 0.0353
    lhl = 0.62
    gap = 0.16
    tot = lhv + gap + lhl
    start = y + 0.30 + max(0, (interior - tot) / 2)
    txt(slide, str(value), x+0.08, start,            w-0.16, lhv+0.22,
        font="Calibri", size=vs, bold=True, color=ac, align=PP_ALIGN.CENTER)
    txt(slide, str(label), x+0.08, start+lhv+gap,    w-0.16, lhl,
        font=T["BF"], size=12, color=tc, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────
# SHARED COVER CONTENT (same text placement for all three cover styles)
# ─────────────────────────────────────────────────────────────────────
def _cover_shared(slide, T, data, offset_x=0.9, max_x=None):
    sw = (max_x if max_x else W - 0.62) - offset_x
    if sw <= 0: return
    u  = safe(data.get("university"))
    f  = " | ".join(filter(None,[safe(data.get("faculty")),safe(data.get("department"))]))
    txt(slide, u, offset_x, 0.26, sw, 1.02,
        font=T["BF"], size=17, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
    if f.strip(" |"):
        txt(slide, f, offset_x, 1.34, sw, 0.80,
            font=T["BF"], size=12, color=T["A"], align=PP_ALIGN.RIGHT)

    bx_w = min(8.2, sw)
    rect(slide, offset_x, 3.04, bx_w, 0.86, T["A"])
    txt(slide, f"مذكرة تخرج — {safe(data.get('level'),'ماستر')}",
        offset_x, 3.04, bx_w, 0.86,
        font=T["BF"], size=14, bold=True, color=T["D"], align=PP_ALIGN.CENTER)

    txt(slide, safe(data.get("titleAr")), offset_x, 4.16, sw, 3.45,
        font=T["BF"], size=20, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)

    if data.get("titleFr"):
        txt(slide, data["titleFr"], offset_x, 7.82, sw, 0.80,
            font="Calibri", size=11, italic=True, color=T["A2"], align=PP_ALIGN.LEFT)

    # Decorative divider
    my = H * 0.555
    ln(slide, offset_x, my, sw, T["A"], 0.07)
    for di in range(3):
        oval(slide, offset_x + di*1.35, my-0.23, 0.44, 0.44, T["A"], 0.5)

    sy = H * 0.790
    ln(slide, offset_x, sy, sw, T["A"], 0.08)
    half = (sw-0.3)/2 - 0.45
    txt(slide, "إعداد الطالب", offset_x, sy+0.26, half, 0.50,
        font=T["BF"], size=11, color=T["A"], align=PP_ALIGN.RIGHT)
    txt(slide, safe(data.get("studentName")), offset_x, sy+0.80, half, 0.86,
        font=T["BF"], size=17, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
    rx = offset_x + half + 0.95; rw = sw - half - 0.95 - 0.3
    txt(slide, "إشراف الأستاذ", rx, sy+0.26, rw, 0.50,
        font=T["BF"], size=11, color=T["A"], align=PP_ALIGN.RIGHT)
    txt(slide, safe(data.get("supervisor")), rx, sy+0.80, rw, 0.86,
        font=T["BF"], size=17, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)

    meta = []
    if data.get("major"): meta.append(f"التخصص: {data['major']}")
    if data.get("year"):  meta.append(f"السنة: {data['year']}")
    if data.get("defenseDate"):
        try:
            d = datetime.date.fromisoformat(data["defenseDate"])
            mo = ["","يناير","فبراير","مارس","أبريل","مايو","يونيو",
                  "يوليو","أغسطس","سبتمبر","أكتوبر","نوفمبر","ديسمبر"]
            meta.append(f"تاريخ المناقشة: {d.day} {mo[d.month]} {d.year}")
        except Exception: meta.append(f"تاريخ المناقشة: {data['defenseDate']}")
    if meta:
        txt(slide, "   ·   ".join(meta), offset_x, sy+1.90, sw, 0.58,
            font=T["BF"], size=11, color=T["TM"], align=PP_ALIGN.RIGHT)
    if data.get("keywords"):
        txt(slide, f"كلمات مفتاحية: {data['keywords']}", offset_x, H-1.15, sw, 0.56,
            font=T["BF"], size=10, color=T["TM"], align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────────
# SHARED QUOTE BOX
# ─────────────────────────────────────────────────────────────────────
def _quote(slide, T, x, y, w, h, text):
    rect(slide, x, y, w, h, T["M"])
    rect(slide, x, y, 0.34, h, T["A"])
    txt(slide, "\u275d", x+0.50, y+0.14, 2.0, 0.98,
        font="Georgia", size=42, color=T["A"], align=PP_ALIGN.RIGHT)
    txt(slide, text, x+0.58, y+1.06, w-0.86, h-1.18,
        font=T["BF"], size=13, color=T["TL"], align=PP_ALIGN.RIGHT)


# ═════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS  (one function per slide type, layout-aware)
# ═════════════════════════════════════════════════════════════════════

def make_cover(prs, data, T):
    slide = blank(prs)
    get_layout(T).cover(slide, T, data)
    return slide


def make_toc(prs, data, T, chapters):
    slide = blank(prs)
    get_layout(T).toc(slide, T, chapters)
    return slide


# ── PROBLEM ──────────────────────────────────────────────────────────
def make_problem(prs, data, T):
    """
    الإشكالية — تخطيط جديد: لوحتان متوازيتان
    · يسار: لمحة الإشكالية + التساؤل الرئيسي (بارز)
    · يمين: التساؤلات الفرعية في بطاقات مرقمة
    """
    slide = blank(prs)
    bg(slide, T["D"])
    oval(slide, W-11, -3, 14, 14, T["M"], 0.38)
    oval(slide, -3, H-8, 10, 10, T["A"], 0.80)
    L = get_layout(T)
    cy0 = L.section_dark(slide, T, "إشكالية البحث والتساؤلات", "Research Problem & Questions")

    problem  = safe(data.get("mainProblem",""))
    main_q   = safe(data.get("mainQuestion",""))
    subs     = [s for s in data.get("subQuestions",[]) if s][:5]
    has_subs = bool(subs)

    if has_subs:
        # تخطيط ثنائي العمود
        lw = W * 0.46 - 0.6   # عرض العمود الأيسر (الإشكالية)
        rw = W - lw - 1.8     # عرض العمود الأيمن (التساؤلات)
        lx, rx = 1.1, lw + 1.5

        # ── لوحة الإشكالية ─────────────────────────────
        _quote(slide, T, lx, cy0, lw, 3.20, problem)

        # ── التساؤل الرئيسي (بارز) ─────────────────────
        if main_q:
            mq_y = cy0 + 3.36
            rect(slide, lx, mq_y, lw, H - mq_y - 0.22, T["A"])
            txt(slide, "؟", lx + 0.08, mq_y + 0.10, 0.72, H - mq_y - 0.42,
                font="Georgia", size=52, bold=True, color=T["D"], align=PP_ALIGN.CENTER)
            txt(slide, main_q, lx + 0.88, mq_y + 0.12, lw - 1.00, H - mq_y - 0.42,
                font=T["BF"], size=13, bold=True, color=T["D"], align=PP_ALIGN.RIGHT)

        # ── التساؤلات الفرعية ──────────────────────────
        avail = H - cy0 - 0.20
        rh = max(1.20, (avail - 0.14 * len(subs)) / len(subs))
        for i, q in enumerate(subs):
            ry = cy0 + i * (rh + 0.14)
            sc = T["SC"][i % len(T["SC"])]
            rect(slide, rx, ry, rw, rh, T["M"] if i%2==0 else T["D"])
            rect(slide, rx, ry, 0.24, rh, sc)
            # رقم دائري
            cx = rx + 0.34; cy = ry + rh/2 - 0.36
            oval(slide, cx, cy, 0.72, 0.72, T["A"], 0.85)
            txt(slide, str(i+1), cx, cy, 0.72, 0.72,
                font="Calibri", size=16, bold=True, color=T["D"], align=PP_ALIGN.CENTER)
            txt(slide, q, rx + 1.22, ry + 0.12, rw - 1.36, rh - 0.24,
                font=T["BF"], size=12, color=T["TL"], align=PP_ALIGN.RIGHT)
    else:
        # تخطيط أحادي — الإشكالية فقط
        _quote(slide, T, 1.1, cy0, W-2.2, H-cy0-0.40, problem)
        if main_q:
            rect(slide, 1.1, H-2.10, W-2.2, 1.88, T["A"])
            txt(slide, main_q, 1.28, H-2.05, W-2.56, 1.72,
                font=T["BF"], size=14, bold=True, color=T["D"], align=PP_ALIGN.RIGHT)
    return slide


# ── OBJECTIVES ───────────────────────────────────────────────────────
def make_objectives(prs, data, T):
    slide = blank(prs)
    if T.get("LAYOUT") in ("bold",):
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "أهداف البحث والفرضيات", "Objectives & Hypotheses")

    objs  = [o for o in data.get("objectives",[])  if o]
    hypos = [h for h in data.get("hypotheses",[])  if h]
    n     = max(len(objs[:6]), len(hypos[:6]), 1)
    cw    = (W-3.08)/2
    avail = H - cy0 - 0.76
    ch    = max(1.26, (avail - 0.16*n) / n)

    # Objectives — left
    rx_lbl = 1.1 if T.get("LAYOUT")!="bold" else 2.0
    rect(slide, rx_lbl, cy0, cw, 0.74, T["D"])
    txt(slide, "🎯  الأهداف", rx_lbl, cy0, cw, 0.74,
        font=T["BF"], size=14, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
    for i, obj in enumerate(objs[:6]):
        oy = cy0+0.76 + i*(ch+0.16)
        sc = T["SC"][i%len(T["SC"])]
        rect(slide, rx_lbl+0.08, oy+0.06, cw, ch, T["CE"])
        rect(slide, rx_lbl,     oy,    cw, ch, T["CB"], line_color=T["CE"])
        rect(slide, rx_lbl,     oy,    0.18, ch, sc)
        nw = 1.70
        rect(slide, rx_lbl, oy, nw, ch, T["D"])
        txt(slide, str(i+1), rx_lbl, oy+ch/2-0.65, nw, 1.3,
            font="Calibri", size=30, bold=True, color=sc, align=PP_ALIGN.CENTER)
        txt(slide, obj, rx_lbl+nw+0.16, oy+0.10, cw-nw-0.28, ch-0.20,
            font=T["BF"], size=12, color=T["TD"], align=PP_ALIGN.RIGHT)

    # Hypotheses — right
    rr = rx_lbl + cw + 0.76
    rect(slide, rr, cy0, cw, 0.74, T["M"])
    txt(slide, "💡  الفرضيات", rr, cy0, cw, 0.74,
        font=T["BF"], size=14, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
    for i, hy in enumerate(hypos[:6]):
        hy_y = cy0+0.76 + i*(ch+0.16)
        sc   = T["SC"][i%len(T["SC"])]
        rect(slide, rr+0.08, hy_y+0.06, cw, ch, T["CE"])
        rect(slide, rr,     hy_y,    cw, ch, T["CB"], line_color=T["CE"])
        rect(slide, rr,     hy_y,    0.18, ch, sc)
        bw = 1.38
        txt(slide, f"H{i+1}", rr+0.26, hy_y+ch*0.08, bw, ch*0.60,
            font="Calibri", size=20, bold=True, color=sc, align=PP_ALIGN.RIGHT)
        txt(slide, hy, rr+bw+0.32, hy_y+0.10, cw-bw-0.56, ch-0.20,
            font=T["BF"], size=12, color=T["TD"], align=PP_ALIGN.RIGHT)
    return slide


# ── IMPORTANCE ───────────────────────────────────────────────────────
def make_importance(prs, data, T):
    slide = blank(prs)
    bg(slide, T["D"])
    oval(slide, W-10, -2, 13, 13, T["M"], 0.50)
    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "أهمية البحث وأسباب اختياره",
                          "Significance & Motivation")
    pairs = [
        ("importance","الأهمية العلمية والعملية","⭐"),
        ("reasons",   "أسباب اختيار الموضوع",    "🔍"),
    ]
    ph = (H-cy0)/2 - 0.26
    for i,(key,lbl,icon) in enumerate(pairs):
        py = cy0 + i*(ph+0.38)
        rect(slide, 1.1, py, W-2.2, ph, T["M"])
        rect(slide, 1.1, py, 0.32, ph, T["A"])
        txt(slide, f"{icon}  {lbl}", 1.56, py+0.15, W-3.28, 0.68,
            font=T["BF"], size=16, bold=True, color=T["A"], align=PP_ALIGN.RIGHT)
        ln(slide, 1.56, py+0.99, W-3.28, T["A"], 0.05)
        txt(slide, safe(data.get(key)), 1.56, py+1.12, W-3.28, ph-1.26,
            font=T["BF"], size=13, color=T["TL"], align=PP_ALIGN.RIGHT)
    return slide


# ── THEORY ───────────────────────────────────────────────────────────
def make_theory(prs, data, T, concepts):
    slide = blank(prs)
    if T.get("LAYOUT") == "bold":
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "الإطار النظري والمفاهيمي",
                            "Theoretical & Conceptual Framework")

    n = min(len(concepts), 6)
    if not n: return slide
    cols = 3 if n >= 3 else n
    rows = math.ceil(n/cols)
    gx, gy = 0.28, 0.24
    x0 = 1.8 if T.get("LAYOUT")=="bold" else 1.0
    aw = W - x0 - 0.8
    cw = (aw - gx*(cols-1))/cols
    avail = H - cy0 - gy*(rows-1)
    ch = avail/rows
    if rows == 1: ch = min(ch, 9.0); grid_y = cy0 + (H-cy0-ch)/2
    else:         grid_y = cy0

    for i, c in enumerate(concepts[:6]):
        col = i%cols; row = i//cols
        cx  = x0 + col*(cw+gx)
        cy  = grid_y + row*(ch+gy)
        sc  = T["SC"][i%len(T["SC"])]
        rect(slide, cx+0.10, cy+0.10, cw, ch, T["CE"])
        rect(slide, cx, cy, cw, ch, T["CB"], line_color=T["CE"])
        rect(slide, cx, cy, cw, 0.78, sc)
        rect(slide, cx, cy, 0.17, ch, T["A"])
        txt(slide, safe(c.get("name")), cx+0.25, cy+0.08, cw-0.42, 0.62,
            font=T["BF"], size=13, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
        txt(slide, safe(c.get("def")), cx+0.25, cy+0.92, cw-0.42, ch-1.06,
            font=T["BF"], size=12, color=T["TD"], align=PP_ALIGN.RIGHT)
    return slide


# ── LITERATURE ───────────────────────────────────────────────────────
def make_literature(prs, data, T, lits):
    slide = blank(prs)
    bg(slide, T["D"])
    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "مراجعة الأدبيات والدراسات السابقة",
                          "Literature Review")

    col_defs = [("الباحث / المؤلف",4.3),("السنة",1.9),
                ("عنوان الدراسة",9.3),("أبرز النتائج",16.1)]
    xs = [1.1]
    for _,cw in col_defs[:-1]: xs.append(xs[-1]+cw+0.10)

    hy,hh = cy0+0.08, 0.84
    for j,((lbl,cw),x) in enumerate(zip(col_defs,xs)):
        sc = T["SC"][j%len(T["SC"])]
        rect(slide, x, hy, cw, hh, sc)
        txt(slide, lbl, x+0.09, hy+0.04, cw-0.18, hh-0.08,
            font=T["BF"], size=12, bold=True, color=T["D"], align=PP_ALIGN.RIGHT)

    n  = min(len(lits), 5)
    rh = max(1.44, (H-hy-hh-0.32)/max(n,1)-0.12)
    for ri, lit in enumerate(lits[:5]):
        ry  = hy+hh+0.10+ri*(rh+0.10)
        bgc = T["M"] if ri%2==0 else T["D"]
        vals = [safe(lit.get("author")), safe(lit.get("year")),
                safe(lit.get("title")), safe(lit.get("findings"))]
        for j,((_, cw),x,val) in enumerate(zip(col_defs,xs,vals)):
            rect(slide, x, ry, cw, rh, bgc)
            rect(slide, x, ry, 0.06, rh, T["SC"][j%len(T["SC"])])
            txt(slide, val, x+0.12, ry+0.08, cw-0.20, rh-0.16,
                font=T["BF"], size=11, color=T["TL"], align=PP_ALIGN.RIGHT)
    return slide


# ── METHODOLOGY ──────────────────────────────────────────────────────
def make_methodology(prs, data, T):
    slide = blank(prs)
    if T.get("LAYOUT") == "bold":
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "المنهجية والأدوات", "Methodology & Tools")
    x0  = 1.8 if T.get("LAYOUT")=="bold" else 1.1

    tests     = [t for t in data.get("statisticalTests",[]) if t]
    has_tests = bool(tests)
    bh = 3.25 if has_tests else (H-cy0)/2 - 0.24
    bw = (W-x0-0.8-0.40)/2

    boxes = [
        ("🔧","المنهج المتبع",  safe(data.get("methodology"))),
        ("📊","مصدر البيانات",  safe(data.get("dataSource"))),
        ("📅","الفترة الزمنية", safe(data.get("timePeriod"))),
        ("💻","برنامج التحليل", safe(data.get("software"))),
    ]
    for i,(icon,lbl,val) in enumerate(boxes):
        bx = x0 + (i%2)*(bw+0.40)
        by = cy0 + (i//2)*(bh+0.28)
        sc = T["SC"][i%len(T["SC"])]
        rect(slide, bx+0.09, by+0.09, bw, bh, T["CE"])
        rect(slide, bx,     by,    bw, bh, T["CB"], line_color=T["CE"])
        rect(slide, bx,     by,    bw, 0.78, sc)
        txt(slide, f"{icon}  {lbl}", bx+0.13, by+0.08, bw-0.26, 0.62,
            font=T["BF"], size=13, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
        txt(slide, val, bx+0.13, by+0.96, bw-0.26, bh-1.10,
            font=T["BF"], size=13, bold=True, color=T["TD"], align=PP_ALIGN.RIGHT)

    if has_tests:
        ty    = cy0 + 2*(bh+0.28)
        avail = H - ty - 0.05
        txt(slide, "الاختبارات الإحصائية المستخدمة", x0, ty, W-x0-0.6, 0.66,
            font=T["BF"], size=14, bold=True, color=T["TD"], align=PP_ALIGN.RIGHT)
        n  = min(len(tests), 5)
        tw = (W-x0-0.6-0.14*(n-1))/n
        th = avail - 0.74
        for i,t in enumerate(tests[:5]):
            tx = x0 + i*(tw+0.14)
            sc = T["SC"][i%len(T["SC"])]
            rect(slide, tx, ty+0.72, tw, th, T["D"])
            rect(slide, tx, ty+0.72, tw, 0.32, sc)
            lines  = max(1, len(t)//15+1)
            text_h = lines*0.56
            text_y = ty+0.72+0.38+max(0,(th-0.38-text_h)/2)
            txt(slide, t, tx+0.08, text_y, tw-0.16, text_h+0.26,
                font=T["BF"], size=11, bold=True, color=T["TL"], align=PP_ALIGN.CENTER)
    return slide


# ── KPI DASHBOARD ────────────────────────────────────────────────────
def make_stats(prs, data, T):
    """لوحة KPI — تصميم داشبورد احترافي مع تأثيرات بصرية"""
    slide = blank(prs)
    bg(slide, T["D"])
    # خلفية هندسية
    oval(slide, W*0.62, -2, W*0.48, H*0.72, T["M"], 0.22)
    rect(slide, 0, 0, W, 0.18, T["A"])

    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "المؤشرات الإحصائية الرئيسية",
                          "Key Performance Indicators — KPI Dashboard")

    stats = [s for s in data.get("stats",[]) if s.get("label") and s.get("value")]
    if not stats: return slide

    n    = min(len(stats), 8)
    cols = min(n, 4)
    rows = math.ceil(n/cols)
    gx, gy = 0.26, 0.26
    cw   = (W - 2.2 - gx*(cols-1)) / cols
    raw_ch = (H - cy0 - gy*(rows-1) - 0.20) / rows
    ch   = min(raw_ch, 3.80)
    tot  = rows*ch + (rows-1)*gy
    gy0  = cy0 + max(0, (H - cy0 - tot - 0.20)/2)

    for i, s in enumerate(stats[:8]):
        col = i % cols
        row = i // cols
        cx  = 1.1 + col*(cw+gx)
        cy  = gy0 + row*(ch+gy)
        sc  = T["SC"][i % len(T["SC"])]

        # ظل البطاقة
        rect(slide, cx+0.08, cy+0.08, cw, ch, T["M"])
        # البطاقة الرئيسية
        rect(slide, cx, cy, cw, ch, T["M"])
        # شريط accent علوي
        rect(slide, cx, cy, cw, 0.16, sc)
        # شريط علوي داخلي
        rect(slide, cx, cy+0.16, cw, 0.06, sc)
        # مستطيل قيمة
        val_h = ch * 0.55
        txt(slide, safe(s["value"]), cx+0.10, cy+0.26, cw-0.20, val_h,
            font="Calibri", size=max(28, min(56, 56 - max(0,len(str(s["value"]))-3)*7)),
            bold=True, color=sc, align=PP_ALIGN.CENTER)
        # خط فاصل
        ln(slide, cx+0.20, cy+0.26+val_h, cw-0.40, sc, 0.030)
        # تسمية
        txt(slide, safe(s["label"]), cx+0.10, cy+0.30+val_h, cw-0.20, ch-0.44-val_h,
            font=T["BF"], size=11, color=T["TM"], align=PP_ALIGN.CENTER)
        # sub نص اختياري
        if s.get("sub"):
            txt(slide, safe(s["sub"]), cx+0.10, cy+ch-0.46, cw-0.20, 0.40,
                font="Calibri", size=9, italic=True, color=T["TM"], align=PP_ALIGN.CENTER)

    return slide


# ── RESULTS ──────────────────────────────────────────────────────────
def make_results(prs, data, T):
    slide = blank(prs)
    if T.get("LAYOUT") == "bold":
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    # ── شريط accent علوي ───────────────────────────────
    rect(slide, 0, 0, W, 0.18, T["A"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "نتائج البحث التفصيلية", "Research Findings")
    x0  = 1.8 if T.get("LAYOUT")=="bold" else 1.1

    results = [r for r in data.get("mainResults",[]) if r]
    n   = min(len(results), 7)
    gap = 0.16
    rh  = max(1.38, (H-cy0-gap*n)/max(n,1))

    for i,res in enumerate(results[:7]):
        ry = cy0 + i*(rh+gap)
        L.result_row(slide, T, x0, ry, W-x0-0.7, rh, i, res)
    return slide


# ── RECOMMENDATIONS ──────────────────────────────────────────────────
def make_recommendations(prs, data, T):
    """التوصيات — بطاقات مرقمة مع accent ملون لكل توصية"""
    recs = [r for r in data.get("recommendations",[]) if r][:6]
    if not recs: return

    slide = blank(prs)
    if T.get("LAYOUT") == "bold":
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    rect(slide, 0, 0, W, 0.18, T["A"])

    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "التوصيات", "Recommendations")
    x0  = 1.8 if T.get("LAYOUT")=="bold" else 1.1

    n    = len(recs)
    cols = 2 if n > 3 else 1
    rows = math.ceil(n/cols)
    gx, gy = 0.30, 0.22
    avail_w = W - x0 - 0.6
    avail_h = H - cy0 - 0.22
    cw = (avail_w - gx*(cols-1)) / cols
    ch = (avail_h - gy*(rows-1)) / rows

    for i, rec in enumerate(recs):
        col = i % cols
        row = i // cols
        cx  = x0 + col*(cw+gx)
        cy  = cy0 + row*(ch+gy)
        sc  = T["SC"][i % len(T["SC"])]

        # ظل
        rect(slide, cx+0.07, cy+0.07, cw, ch, T["CE"])
        # بطاقة
        rect(slide, cx, cy, cw, ch, T["CB"], line_color=T["CE"])
        # شريط accent يمين
        rect(slide, cx, cy, 0.26, ch, sc)
        # رقم ترتيبي
        oval(slide, cx+0.36, cy+ch/2-0.44, 0.88, 0.88, sc, 0.90)
        txt(slide, f"{i+1:02d}", cx+0.36, cy+ch/2-0.44, 0.88, 0.88,
            font="Calibri", size=20, bold=True, color=T["TL"] if T.get("LAYOUT")!="minimal" else T["D"],
            align=PP_ALIGN.CENTER)
        # نص
        txt(slide, safe(rec), cx+1.36, cy+0.12, cw-1.52, ch-0.24,
            font=T["BF"], size=12, color=T["TD"], align=PP_ALIGN.RIGHT)
        # تزيين: خط accent أسفل
        ln(slide, cx+1.36, cy+ch-0.22, cw-1.52, sc, 0.030)
    return slide


def make_future(prs, data, T):
    slide = blank(prs)
    if T.get("LAYOUT") == "bold":
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "آفاق وامتدادات البحث",
                            "Future Research Perspectives")
    x0  = 1.8 if T.get("LAYOUT")=="bold" else 1.1

    futures = [f for f in data.get("futureWork",[]) if f]
    n = min(len(futures), 5)
    if not n: return slide

    tlx  = W - 3.08
    rect(slide, tlx-0.06, cy0, 0.12, H-cy0-0.36, T["A"])
    gap = 0.24
    fh  = (H-cy0-0.36-gap*n)/n

    for i,fw in enumerate(futures[:5]):
        fy  = cy0 + i*(fh+gap)
        ncy = fy + fh/2 - 0.43
        sc  = T["SC"][i%len(T["SC"])]
        oval(slide, tlx-0.45, ncy, 0.90, 0.90, T["D"])
        oval(slide, tlx-0.35, ncy+0.10, 0.70, 0.70, sc)
        txt(slide, str(i+1), tlx-0.35, ncy+0.10, 0.70, 0.70,
            font="Calibri", size=12, bold=True, color=T["D"], align=PP_ALIGN.CENTER)
        rect(slide, tlx-1.68, ncy+0.37, 1.26, 0.10, T["A"])
        cw = tlx - x0 - 1.68
        rect(slide, x0, fy+0.06, cw, fh, T["CB"], line_color=T["CE"])
        rect(slide, x0, fy+0.06, 0.19, fh, sc)
        txt(slide, f"آفق بحثي {i+1}", x0+0.28, fy+0.10, 4.4, 0.48,
            font=T["BF"], size=10, bold=True, color=sc, align=PP_ALIGN.RIGHT)
        txt(slide, fw, x0+0.28, fy+0.60, cw-0.50, fh-0.70,
            font=T["BF"], size=12, color=T["TD"], align=PP_ALIGN.RIGHT)
    return slide


# ── CONCLUSION ───────────────────────────────────────────────────────
def make_conclusion(prs, data, T):
    """الخاتمة — تصميم سينمائي: خلفية داكنة مع اقتباس بارز"""
    slide = blank(prs)
    bg(slide, T["D"])
    oval(slide, W*0.60, -2, W*0.50, H*0.75, T["M"], 0.30)
    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "الخاتمة والاستنتاجات", "Conclusion")

    conclusion = safe(data.get("generalConclusion",""))
    recs  = [r for r in data.get("recommendations",[]) if r][:4]

    # لوحة الاقتباس الرئيسية
    qh = 3.80 if recs else H - cy0 - 0.35
    rect(slide, 1.1, cy0, W-2.2, qh, T["M"])
    rect(slide, 1.1, cy0, 0.38, qh, T["A"])
    txt(slide, "“", 1.6, cy0+0.08, 2.0, 1.20,
        font="Georgia", size=72, bold=True, color=T["A"], align=PP_ALIGN.RIGHT)
    txt(slide, conclusion, 1.58, cy0+1.10, W-2.9, qh-1.30,
        font=T["BF"], size=14, color=T["TL"], align=PP_ALIGN.RIGHT)
    txt(slide, "”", W-2.8, cy0+qh-1.0, 1.50, 1.0,
        font="Georgia", size=72, bold=True, color=T["A"], align=PP_ALIGN.LEFT)

    # أبرز التوصيات (اختصار)
    if recs:
        ry = cy0 + qh + 0.22
        avail = H - ry - 0.18
        rh = avail / len(recs)
        rw = (W - 2.4 - 0.20*(len(recs)-1)) / len(recs)
        for i, rec in enumerate(recs):
            rx = 1.1 + i*(rw+0.20)
            sc = T["SC"][i % len(T["SC"])]
            rect(slide, rx, ry, rw, rh, T["M"])
            rect(slide, rx, ry, rw, 0.055, sc)
            rect(slide, rx, ry, 0.14, rh, sc)
            txt(slide, f"{i+1:02d}", rx+0.22, ry+0.06, rw-0.34, 0.52,
                font="Calibri", size=20, bold=True, color=sc, align=PP_ALIGN.RIGHT)
            txt(slide, safe(rec), rx+0.22, ry+0.62, rw-0.34, rh-0.76,
                font=T["BF"], size=11, color=T["TL"], align=PP_ALIGN.RIGHT)
    return slide


def make_final(prs, data, T):
    """شريحة الشكر — تصميم سينمائي احترافي"""
    slide = blank(prs)
    bg(slide, T["D"])
    # ── خلفية هندسية ─────────────────────────────────
    oval(slide, W*0.55, -4, W*0.60, H*0.90, T["M"], 0.28)
    oval(slide, -5, H*0.50, W*0.70, W*0.70, T["A"], 0.55)
    rect(slide, 0, 0, W, 0.55, T["A"])
    rect(slide, 0, H-0.55, W, 0.55, T["A"])

    # ── نص الشكر الرئيسي ─────────────────────────────
    txt(slide, "شكراً لحسن استماعكم", 1.0, H*0.24, W-2.0, 1.80,
        font=T["HF"], size=48, bold=True, color=T["TL"], align=PP_ALIGN.CENTER)
    txt(slide, "Merci pour votre attention", 1.0, H*0.24+1.90, W-2.0, 0.80,
        font="Calibri", size=22, italic=True, color=T["A"], align=PP_ALIGN.CENTER)
    txt(slide, "Thank you", 1.0, H*0.24+2.76, W-2.0, 0.60,
        font="Calibri", size=16, italic=True, color=T["TM"], align=PP_ALIGN.CENTER)

    # ── فاصل ─────────────────────────────────────────
    ln(slide, W*0.25, H*0.73, W*0.50, T["A"], 0.06)

    # ── معلومات الطالب ───────────────────────────────
    student = safe(data.get("studentName",""))
    supervisor = safe(data.get("supervisor",""))
    year = safe(data.get("year",""))
    info_parts = []
    if student:    info_parts.append(f"إعداد: {student}")
    if supervisor: info_parts.append(f"إشراف: {supervisor}")
    if year:       info_parts.append(year)
    info = "   ·   ".join(info_parts)
    txt(slide, info, 1.0, H*0.77, W-2.0, 0.60,
        font=T["BF"], size=13, bold=False, color=T["TM"], align=PP_ALIGN.CENTER)

    # ── اللقب الأكاديمي ───────────────────────────────
    major = safe(data.get("major",""))
    univ  = safe(data.get("university",""))
    if major or univ:
        acad = " — ".join(filter(None, [major, univ]))
        txt(slide, acad, 1.0, H*0.83, W-2.0, 0.48,
            font=T["BF"], size=11, italic=True, color=T["A"], align=PP_ALIGN.CENTER)
    return slide






def make_importance_v2(prs, data, T):
    """أهمية الدراسة — بطاقات بصرية متدرجة مع أيقونات"""
    items_list = [x for x in data.get("importance",[]) if x]
    reasons    = safe(data.get("reasons",""))
    all_items  = items_list[:]
    if reasons: all_items.append(reasons)
    if not all_items: return

    slide = blank(prs)
    bg(slide, T["D"])
    oval(slide, -3, H*0.40, W*0.55, W*0.55, T["M"], 0.28)
    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "أهمية الدراسة", "Research Significance")

    n    = min(len(all_items), 6)
    cols = 2 if n > 3 else 1
    rows = math.ceil(n/cols)
    gx, gy = 0.32, 0.20
    avail_w = W - 2.2
    avail_h = H - cy0 - 0.22
    cw = (avail_w - gx*(cols-1)) / cols
    ch = (avail_h - gy*(rows-1)) / rows

    icons = ["🔬","💡","📊","🎯","🌐","⚡"]
    for i, item in enumerate(all_items[:6]):
        col = i % cols
        row = i // cols
        cx  = 1.1 + col*(cw+gx)
        cy  = cy0 + row*(ch+gy)
        sc  = T["SC"][i % len(T["SC"])]

        # بطاقة مع تدرج
        rect(slide, cx, cy, cw, ch, T["M"])
        rect(slide, cx, cy, 0.28, ch, sc)           # شريط جانبي
        rect(slide, cx, cy, cw, 0.055, sc)           # شريط علوي

        # أيقونة + رقم
        icon = icons[i % len(icons)]
        txt(slide, icon, cx+0.36, cy+0.08, 0.80, 0.70,
            font="Segoe UI Emoji", size=22, align=PP_ALIGN.CENTER)
        oval(slide, cx+cw-0.72, cy+0.10, 0.58, 0.58, T["D"], 0.85)
        txt(slide, f"{i+1:02d}", cx+cw-0.72, cy+0.10, 0.58, 0.58,
            font="Calibri", size=14, bold=True, color=sc, align=PP_ALIGN.CENTER)

        # خط فاصل
        ln(slide, cx+0.36, cy+0.82, cw-0.50, sc, 0.025)

        # نص
        txt(slide, safe(item), cx+0.36, cy+0.90, cw-0.52, ch-1.02,
            font=T["BF"], size=12, color=T["TL"], align=PP_ALIGN.RIGHT)
    return slide


def make_cover(prs, data, T):
    slide = blank(prs)
    get_layout(T).cover(slide, T, data)
    return slide


def make_toc(prs, data, T, chapters):
    slide = blank(prs)
    get_layout(T).toc(slide, T, chapters)
    return slide


# ── PROBLEM ──────────────────────────────────────────────────────────
def make_problem(prs, data, T):
    """
    الإشكالية — تخطيط جديد: لوحتان متوازيتان
    · يسار: لمحة الإشكالية + التساؤل الرئيسي (بارز)
    · يمين: التساؤلات الفرعية في بطاقات مرقمة
    """
    slide = blank(prs)
    bg(slide, T["D"])
    oval(slide, W-11, -3, 14, 14, T["M"], 0.38)
    oval(slide, -3, H-8, 10, 10, T["A"], 0.80)
    L = get_layout(T)
    cy0 = L.section_dark(slide, T, "إشكالية البحث والتساؤلات", "Research Problem & Questions")

    problem  = safe(data.get("mainProblem",""))
    main_q   = safe(data.get("mainQuestion",""))
    subs     = [s for s in data.get("subQuestions",[]) if s][:5]
    has_subs = bool(subs)

    if has_subs:
        # تخطيط ثنائي العمود
        lw = W * 0.46 - 0.6   # عرض العمود الأيسر (الإشكالية)
        rw = W - lw - 1.8     # عرض العمود الأيمن (التساؤلات)
        lx, rx = 1.1, lw + 1.5

        # ── لوحة الإشكالية ─────────────────────────────
        _quote(slide, T, lx, cy0, lw, 3.20, problem)

        # ── التساؤل الرئيسي (بارز) ─────────────────────
        if main_q:
            mq_y = cy0 + 3.36
            rect(slide, lx, mq_y, lw, H - mq_y - 0.22, T["A"])
            txt(slide, "؟", lx + 0.08, mq_y + 0.10, 0.72, H - mq_y - 0.42,
                font="Georgia", size=52, bold=True, color=T["D"], align=PP_ALIGN.CENTER)
            txt(slide, main_q, lx + 0.88, mq_y + 0.12, lw - 1.00, H - mq_y - 0.42,
                font=T["BF"], size=13, bold=True, color=T["D"], align=PP_ALIGN.RIGHT)

        # ── التساؤلات الفرعية ──────────────────────────
        avail = H - cy0 - 0.20
        rh = max(1.20, (avail - 0.14 * len(subs)) / len(subs))
        for i, q in enumerate(subs):
            ry = cy0 + i * (rh + 0.14)
            sc = T["SC"][i % len(T["SC"])]
            rect(slide, rx, ry, rw, rh, T["M"] if i%2==0 else T["D"])
            rect(slide, rx, ry, 0.24, rh, sc)
            # رقم دائري
            cx = rx + 0.34; cy = ry + rh/2 - 0.36
            oval(slide, cx, cy, 0.72, 0.72, T["A"], 0.85)
            txt(slide, str(i+1), cx, cy, 0.72, 0.72,
                font="Calibri", size=16, bold=True, color=T["D"], align=PP_ALIGN.CENTER)
            txt(slide, q, rx + 1.22, ry + 0.12, rw - 1.36, rh - 0.24,
                font=T["BF"], size=12, color=T["TL"], align=PP_ALIGN.RIGHT)
    else:
        # تخطيط أحادي — الإشكالية فقط
        _quote(slide, T, 1.1, cy0, W-2.2, H-cy0-0.40, problem)
        if main_q:
            rect(slide, 1.1, H-2.10, W-2.2, 1.88, T["A"])
            txt(slide, main_q, 1.28, H-2.05, W-2.56, 1.72,
                font=T["BF"], size=14, bold=True, color=T["D"], align=PP_ALIGN.RIGHT)
    return slide


# ── OBJECTIVES ───────────────────────────────────────────────────────
def make_objectives(prs, data, T):
    slide = blank(prs)
    if T.get("LAYOUT") in ("bold",):
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "أهداف البحث والفرضيات", "Objectives & Hypotheses")

    objs  = [o for o in data.get("objectives",[])  if o]
    hypos = [h for h in data.get("hypotheses",[])  if h]
    n     = max(len(objs[:6]), len(hypos[:6]), 1)
    cw    = (W-3.08)/2
    avail = H - cy0 - 0.76
    ch    = max(1.26, (avail - 0.16*n) / n)

    # Objectives — left
    rx_lbl = 1.1 if T.get("LAYOUT")!="bold" else 2.0
    rect(slide, rx_lbl, cy0, cw, 0.74, T["D"])
    txt(slide, "🎯  الأهداف", rx_lbl, cy0, cw, 0.74,
        font=T["BF"], size=14, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
    for i, obj in enumerate(objs[:6]):
        oy = cy0+0.76 + i*(ch+0.16)
        sc = T["SC"][i%len(T["SC"])]
        rect(slide, rx_lbl+0.08, oy+0.06, cw, ch, T["CE"])
        rect(slide, rx_lbl,     oy,    cw, ch, T["CB"], line_color=T["CE"])
        rect(slide, rx_lbl,     oy,    0.18, ch, sc)
        nw = 1.70
        rect(slide, rx_lbl, oy, nw, ch, T["D"])
        txt(slide, str(i+1), rx_lbl, oy+ch/2-0.65, nw, 1.3,
            font="Calibri", size=30, bold=True, color=sc, align=PP_ALIGN.CENTER)
        txt(slide, obj, rx_lbl+nw+0.16, oy+0.10, cw-nw-0.28, ch-0.20,
            font=T["BF"], size=12, color=T["TD"], align=PP_ALIGN.RIGHT)

    # Hypotheses — right
    rr = rx_lbl + cw + 0.76
    rect(slide, rr, cy0, cw, 0.74, T["M"])
    txt(slide, "💡  الفرضيات", rr, cy0, cw, 0.74,
        font=T["BF"], size=14, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
    for i, hy in enumerate(hypos[:6]):
        hy_y = cy0+0.76 + i*(ch+0.16)
        sc   = T["SC"][i%len(T["SC"])]
        rect(slide, rr+0.08, hy_y+0.06, cw, ch, T["CE"])
        rect(slide, rr,     hy_y,    cw, ch, T["CB"], line_color=T["CE"])
        rect(slide, rr,     hy_y,    0.18, ch, sc)
        bw = 1.38
        txt(slide, f"H{i+1}", rr+0.26, hy_y+ch*0.08, bw, ch*0.60,
            font="Calibri", size=20, bold=True, color=sc, align=PP_ALIGN.RIGHT)
        txt(slide, hy, rr+bw+0.32, hy_y+0.10, cw-bw-0.56, ch-0.20,
            font=T["BF"], size=12, color=T["TD"], align=PP_ALIGN.RIGHT)
    return slide


# ── IMPORTANCE ───────────────────────────────────────────────────────
def make_importance(prs, data, T):
    slide = blank(prs)
    bg(slide, T["D"])
    oval(slide, W-10, -2, 13, 13, T["M"], 0.50)
    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "أهمية البحث وأسباب اختياره",
                          "Significance & Motivation")
    pairs = [
        ("importance","الأهمية العلمية والعملية","⭐"),
        ("reasons",   "أسباب اختيار الموضوع",    "🔍"),
    ]
    ph = (H-cy0)/2 - 0.26
    for i,(key,lbl,icon) in enumerate(pairs):
        py = cy0 + i*(ph+0.38)
        rect(slide, 1.1, py, W-2.2, ph, T["M"])
        rect(slide, 1.1, py, 0.32, ph, T["A"])
        txt(slide, f"{icon}  {lbl}", 1.56, py+0.15, W-3.28, 0.68,
            font=T["BF"], size=16, bold=True, color=T["A"], align=PP_ALIGN.RIGHT)
        ln(slide, 1.56, py+0.99, W-3.28, T["A"], 0.05)
        txt(slide, safe(data.get(key)), 1.56, py+1.12, W-3.28, ph-1.26,
            font=T["BF"], size=13, color=T["TL"], align=PP_ALIGN.RIGHT)
    return slide


# ── THEORY ───────────────────────────────────────────────────────────
def make_theory(prs, data, T, concepts):
    slide = blank(prs)
    if T.get("LAYOUT") == "bold":
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "الإطار النظري والمفاهيمي",
                            "Theoretical & Conceptual Framework")

    n = min(len(concepts), 6)
    if not n: return slide
    cols = 3 if n >= 3 else n
    rows = math.ceil(n/cols)
    gx, gy = 0.28, 0.24
    x0 = 1.8 if T.get("LAYOUT")=="bold" else 1.0
    aw = W - x0 - 0.8
    cw = (aw - gx*(cols-1))/cols
    avail = H - cy0 - gy*(rows-1)
    ch = avail/rows
    if rows == 1: ch = min(ch, 9.0); grid_y = cy0 + (H-cy0-ch)/2
    else:         grid_y = cy0

    for i, c in enumerate(concepts[:6]):
        col = i%cols; row = i//cols
        cx  = x0 + col*(cw+gx)
        cy  = grid_y + row*(ch+gy)
        sc  = T["SC"][i%len(T["SC"])]
        rect(slide, cx+0.10, cy+0.10, cw, ch, T["CE"])
        rect(slide, cx, cy, cw, ch, T["CB"], line_color=T["CE"])
        rect(slide, cx, cy, cw, 0.78, sc)
        rect(slide, cx, cy, 0.17, ch, T["A"])
        txt(slide, safe(c.get("name")), cx+0.25, cy+0.08, cw-0.42, 0.62,
            font=T["BF"], size=13, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
        txt(slide, safe(c.get("def")), cx+0.25, cy+0.92, cw-0.42, ch-1.06,
            font=T["BF"], size=12, color=T["TD"], align=PP_ALIGN.RIGHT)
    return slide


# ── LITERATURE ───────────────────────────────────────────────────────
def make_literature(prs, data, T, lits):
    slide = blank(prs)
    bg(slide, T["D"])
    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "مراجعة الأدبيات والدراسات السابقة",
                          "Literature Review")

    col_defs = [("الباحث / المؤلف",4.3),("السنة",1.9),
                ("عنوان الدراسة",9.3),("أبرز النتائج",16.1)]
    xs = [1.1]
    for _,cw in col_defs[:-1]: xs.append(xs[-1]+cw+0.10)

    hy,hh = cy0+0.08, 0.84
    for j,((lbl,cw),x) in enumerate(zip(col_defs,xs)):
        sc = T["SC"][j%len(T["SC"])]
        rect(slide, x, hy, cw, hh, sc)
        txt(slide, lbl, x+0.09, hy+0.04, cw-0.18, hh-0.08,
            font=T["BF"], size=12, bold=True, color=T["D"], align=PP_ALIGN.RIGHT)

    n  = min(len(lits), 5)
    rh = max(1.44, (H-hy-hh-0.32)/max(n,1)-0.12)
    for ri, lit in enumerate(lits[:5]):
        ry  = hy+hh+0.10+ri*(rh+0.10)
        bgc = T["M"] if ri%2==0 else T["D"]
        vals = [safe(lit.get("author")), safe(lit.get("year")),
                safe(lit.get("title")), safe(lit.get("findings"))]
        for j,((_, cw),x,val) in enumerate(zip(col_defs,xs,vals)):
            rect(slide, x, ry, cw, rh, bgc)
            rect(slide, x, ry, 0.06, rh, T["SC"][j%len(T["SC"])])
            txt(slide, val, x+0.12, ry+0.08, cw-0.20, rh-0.16,
                font=T["BF"], size=11, color=T["TL"], align=PP_ALIGN.RIGHT)
    return slide


# ── METHODOLOGY ──────────────────────────────────────────────────────
def make_methodology(prs, data, T):
    slide = blank(prs)
    if T.get("LAYOUT") == "bold":
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "المنهجية والأدوات", "Methodology & Tools")
    x0  = 1.8 if T.get("LAYOUT")=="bold" else 1.1

    tests     = [t for t in data.get("statisticalTests",[]) if t]
    has_tests = bool(tests)
    bh = 3.25 if has_tests else (H-cy0)/2 - 0.24
    bw = (W-x0-0.8-0.40)/2

    boxes = [
        ("🔧","المنهج المتبع",  safe(data.get("methodology"))),
        ("📊","مصدر البيانات",  safe(data.get("dataSource"))),
        ("📅","الفترة الزمنية", safe(data.get("timePeriod"))),
        ("💻","برنامج التحليل", safe(data.get("software"))),
    ]
    for i,(icon,lbl,val) in enumerate(boxes):
        bx = x0 + (i%2)*(bw+0.40)
        by = cy0 + (i//2)*(bh+0.28)
        sc = T["SC"][i%len(T["SC"])]
        rect(slide, bx+0.09, by+0.09, bw, bh, T["CE"])
        rect(slide, bx,     by,    bw, bh, T["CB"], line_color=T["CE"])
        rect(slide, bx,     by,    bw, 0.78, sc)
        txt(slide, f"{icon}  {lbl}", bx+0.13, by+0.08, bw-0.26, 0.62,
            font=T["BF"], size=13, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
        txt(slide, val, bx+0.13, by+0.96, bw-0.26, bh-1.10,
            font=T["BF"], size=13, bold=True, color=T["TD"], align=PP_ALIGN.RIGHT)

    if has_tests:
        ty    = cy0 + 2*(bh+0.28)
        avail = H - ty - 0.05
        txt(slide, "الاختبارات الإحصائية المستخدمة", x0, ty, W-x0-0.6, 0.66,
            font=T["BF"], size=14, bold=True, color=T["TD"], align=PP_ALIGN.RIGHT)
        n  = min(len(tests), 5)
        tw = (W-x0-0.6-0.14*(n-1))/n
        th = avail - 0.74
        for i,t in enumerate(tests[:5]):
            tx = x0 + i*(tw+0.14)
            sc = T["SC"][i%len(T["SC"])]
            rect(slide, tx, ty+0.72, tw, th, T["D"])
            rect(slide, tx, ty+0.72, tw, 0.32, sc)
            lines  = max(1, len(t)//15+1)
            text_h = lines*0.56
            text_y = ty+0.72+0.38+max(0,(th-0.38-text_h)/2)
            txt(slide, t, tx+0.08, text_y, tw-0.16, text_h+0.26,
                font=T["BF"], size=11, bold=True, color=T["TL"], align=PP_ALIGN.CENTER)
    return slide


# ── KPI DASHBOARD ────────────────────────────────────────────────────
def make_stats(prs, data, T):
    slide = blank(prs)
    bg(slide, T["D"])
    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "النتائج الكمية والإحصائية",
                          "Key Statistical Results — KPI Dashboard")

    stats = [s for s in data.get("stats",[]) if s.get("label") and s.get("value")]
    if not stats: return slide

    n    = min(len(stats), 8)
    cols = min(n, 4); rows = math.ceil(n/cols)
    gx, gy = 0.24, 0.32
    cw = (W-2.2-gx*(cols-1))/cols
    raw_ch = (H-cy0-gy*(rows-1))/rows
    ch     = min(raw_ch, 7.0)
    tot    = rows*ch + (rows-1)*gy
    grid_y = cy0 + max(0, (H-cy0-tot)/2)

    for i,s in enumerate(stats[:8]):
        col = i%cols; row = i//cols
        cx  = 1.1 + col*(cw+gx)
        cy  = grid_y + row*(ch+gy)
        sc  = T["SC"][i%len(T["SC"])]
        L.kpi_card(slide, T, cx, cy, cw, ch, s["value"], s["label"], sc)
    return slide


# ── RESULTS ──────────────────────────────────────────────────────────
def make_results(prs, data, T):
    slide = blank(prs)
    if T.get("LAYOUT") == "bold":
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    # ── شريط accent علوي ───────────────────────────────
    rect(slide, 0, 0, W, 0.18, T["A"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "نتائج البحث التفصيلية", "Research Findings")
    x0  = 1.8 if T.get("LAYOUT")=="bold" else 1.1

    results = [r for r in data.get("mainResults",[]) if r]
    n   = min(len(results), 7)
    gap = 0.16
    rh  = max(1.38, (H-cy0-gap*n)/max(n,1))

    for i,res in enumerate(results[:7]):
        ry = cy0 + i*(rh+gap)
        L.result_row(slide, T, x0, ry, W-x0-0.7, rh, i, res)
    return slide


# ── RECOMMENDATIONS ──────────────────────────────────────────────────
def make_recommendations(prs, data, T):
    slide = blank(prs)
    bg(slide, T["D"])
    oval(slide, W-12, -4, 16, 16, T["M"], 0.45)
    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "التوصيات", "Recommendations")

    recs = [r for r in data.get("recommendations",[]) if r]
    n    = min(len(recs), 6)
    if not n: return slide

    cols = 2; cw = (W-3.06)/cols
    gapx = 0.72; gapy = 0.28
    rows = math.ceil(n/cols)
    ch   = (H-cy0-gapy*rows)/rows

    for i,rec in enumerate(recs[:6]):
        col = i%cols; row = i//cols
        cx  = 1.1 + col*(cw+gapx)
        cy  = cy0 + row*(ch+gapy)
        sc  = T["SC"][i%len(T["SC"])]
        rect(slide, cx, cy, cw, ch, T["M"])
        rect(slide, cx, cy, 0.25, ch, sc)
        txt(slide, f"{i+1:02d}", cx+0.34, cy+0.10, 1.44, ch*0.40,
            font="Calibri", size=26, bold=True, color=sc, align=PP_ALIGN.LEFT)
        ln(slide, cx+0.34, cy+ch*0.42, cw-0.50, sc, 0.05)
        txt(slide, rec, cx+0.34, cy+ch*0.48, cw-0.50, ch*0.48,
            font=T["BF"], size=12, color=T["TL"], align=PP_ALIGN.RIGHT)
    return slide


# ── FUTURE PERSPECTIVES ──────────────────────────────────────────────
def make_future(prs, data, T):
    slide = blank(prs)
    if T.get("LAYOUT") == "bold":
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "آفاق وامتدادات البحث",
                            "Future Research Perspectives")
    x0  = 1.8 if T.get("LAYOUT")=="bold" else 1.1

    futures = [f for f in data.get("futureWork",[]) if f]
    n = min(len(futures), 5)
    if not n: return slide

    tlx  = W - 3.08
    rect(slide, tlx-0.06, cy0, 0.12, H-cy0-0.36, T["A"])
    gap = 0.24
    fh  = (H-cy0-0.36-gap*n)/n

    for i,fw in enumerate(futures[:5]):
        fy  = cy0 + i*(fh+gap)
        ncy = fy + fh/2 - 0.43
        sc  = T["SC"][i%len(T["SC"])]
        oval(slide, tlx-0.45, ncy, 0.90, 0.90, T["D"])
        oval(slide, tlx-0.35, ncy+0.10, 0.70, 0.70, sc)
        txt(slide, str(i+1), tlx-0.35, ncy+0.10, 0.70, 0.70,
            font="Calibri", size=12, bold=True, color=T["D"], align=PP_ALIGN.CENTER)
        rect(slide, tlx-1.68, ncy+0.37, 1.26, 0.10, T["A"])
        cw = tlx - x0 - 1.68
        rect(slide, x0, fy+0.06, cw, fh, T["CB"], line_color=T["CE"])
        rect(slide, x0, fy+0.06, 0.19, fh, sc)
        txt(slide, f"آفق بحثي {i+1}", x0+0.28, fy+0.10, 4.4, 0.48,
            font=T["BF"], size=10, bold=True, color=sc, align=PP_ALIGN.RIGHT)
        txt(slide, fw, x0+0.28, fy+0.60, cw-0.50, fh-0.70,
            font=T["BF"], size=12, color=T["TD"], align=PP_ALIGN.RIGHT)
    return slide


# ── CONCLUSION ───────────────────────────────────────────────────────
def make_conclusion(prs, data, T):
    """الخاتمة — تصميم سينمائي: خلفية داكنة مع اقتباس بارز"""
    slide = blank(prs)
    bg(slide, T["D"])
    oval(slide, W*0.60, -2, W*0.50, H*0.75, T["M"], 0.30)
    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "الخاتمة والاستنتاجات", "Conclusion")

    conclusion = safe(data.get("generalConclusion",""))
    recs  = [r for r in data.get("recommendations",[]) if r][:4]

    # لوحة الاقتباس الرئيسية
    qh = 3.80 if recs else H - cy0 - 0.35
    rect(slide, 1.1, cy0, W-2.2, qh, T["M"])
    rect(slide, 1.1, cy0, 0.38, qh, T["A"])
    txt(slide, "“", 1.6, cy0+0.08, 2.0, 1.20,
        font="Georgia", size=72, bold=True, color=T["A"], align=PP_ALIGN.RIGHT)
    txt(slide, conclusion, 1.58, cy0+1.10, W-2.9, qh-1.30,
        font=T["BF"], size=14, color=T["TL"], align=PP_ALIGN.RIGHT)
    txt(slide, "”", W-2.8, cy0+qh-1.0, 1.50, 1.0,
        font="Georgia", size=72, bold=True, color=T["A"], align=PP_ALIGN.LEFT)

    # أبرز التوصيات (اختصار)
    if recs:
        ry = cy0 + qh + 0.22
        avail = H - ry - 0.18
        rh = avail / len(recs)
        rw = (W - 2.4 - 0.20*(len(recs)-1)) / len(recs)
        for i, rec in enumerate(recs):
            rx = 1.1 + i*(rw+0.20)
            sc = T["SC"][i % len(T["SC"])]
            rect(slide, rx, ry, rw, rh, T["M"])
            rect(slide, rx, ry, rw, 0.055, sc)
            rect(slide, rx, ry, 0.14, rh, sc)
            txt(slide, f"{i+1:02d}", rx+0.22, ry+0.06, rw-0.34, 0.52,
                font="Calibri", size=20, bold=True, color=sc, align=PP_ALIGN.RIGHT)
            txt(slide, safe(rec), rx+0.22, ry+0.62, rw-0.34, rh-0.76,
                font=T["BF"], size=11, color=T["TL"], align=PP_ALIGN.RIGHT)
    return slide




def make_final(prs, data, T):
    """شريحة الشكر — تصميم سينمائي احترافي"""
    slide = blank(prs)
    bg(slide, T["D"])
    # ── خلفية هندسية ─────────────────────────────────
    oval(slide, W*0.55, -4, W*0.60, H*0.90, T["M"], 0.28)
    oval(slide, -5, H*0.50, W*0.70, W*0.70, T["A"], 0.55)
    rect(slide, 0, 0, W, 0.55, T["A"])
    rect(slide, 0, H-0.55, W, 0.55, T["A"])

    # ── نص الشكر الرئيسي ─────────────────────────────
    txt(slide, "شكراً لحسن استماعكم", 1.0, H*0.24, W-2.0, 1.80,
        font=T["HF"], size=48, bold=True, color=T["TL"], align=PP_ALIGN.CENTER)
    txt(slide, "Merci pour votre attention", 1.0, H*0.24+1.90, W-2.0, 0.80,
        font="Calibri", size=22, italic=True, color=T["A"], align=PP_ALIGN.CENTER)
    txt(slide, "Thank you", 1.0, H*0.24+2.76, W-2.0, 0.60,
        font="Calibri", size=16, italic=True, color=T["TM"], align=PP_ALIGN.CENTER)

    # ── فاصل ─────────────────────────────────────────
    ln(slide, W*0.25, H*0.73, W*0.50, T["A"], 0.06)

    # ── معلومات الطالب ───────────────────────────────
    student = safe(data.get("studentName",""))
    supervisor = safe(data.get("supervisor",""))
    year = safe(data.get("year",""))
    info_parts = []
    if student:    info_parts.append(f"إعداد: {student}")
    if supervisor: info_parts.append(f"إشراف: {supervisor}")
    if year:       info_parts.append(year)
    info = "   ·   ".join(info_parts)
    txt(slide, info, 1.0, H*0.77, W-2.0, 0.60,
        font=T["BF"], size=13, bold=False, color=T["TM"], align=PP_ALIGN.CENTER)

    # ── اللقب الأكاديمي ───────────────────────────────
    major = safe(data.get("major",""))
    univ  = safe(data.get("university",""))
    if major or univ:
        acad = " — ".join(filter(None, [major, univ]))
        txt(slide, acad, 1.0, H*0.83, W-2.0, 0.48,
            font=T["BF"], size=11, italic=True, color=T["A"], align=PP_ALIGN.CENTER)
    return slide







# ─────────────────────────────────────────────────────────────────────
# make_intro — شريحة المقدمة
# ─────────────────────────────────────────────────────────────────────
def make_intro(prs, data, T):
    slide = blank(prs)
    bg(slide, T["D"])
    oval(slide, W*0.62, -2, W*0.48, H*0.72, T["M"], 0.28)
    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "المقدمة", "Introduction")
    overview = safe(data.get("introOverview", ""))
    approach = safe(data.get("introApproach", ""))
    if overview:
        txt(slide, "لمحة عامة عن الموضوع", 1.1, cy0, W-2.2, 0.56,
            font=T["HF"], size=14, bold=True, color=T["A"], align=PP_ALIGN.RIGHT)
        txt(slide, overview, 1.1, cy0+0.60, W-2.2,
            H - cy0 - (1.65 if approach else 0.40) - 0.60,
            font=T["BF"], size=13, color=T["TL"], align=PP_ALIGN.RIGHT)
    if approach:
        ay = H - 1.60
        rect(slide, 1.1, ay, W-2.2, 1.44, T["M"])
        rect(slide, 1.1, ay, 0.22, 1.44, T["A"])
        txt(slide, "المقاربة النظرية", 1.42, ay+0.10, W-2.8, 0.42,
            font=T["BF"], size=13, bold=True, color=T["A"], align=PP_ALIGN.RIGHT)
        txt(slide, approach, 1.42, ay+0.54, W-2.8, 0.80,
            font=T["BF"], size=12, color=T["TL"], align=PP_ALIGN.RIGHT)
    return slide


# ─────────────────────────────────────────────────────────────────────
# make_plan — خطة الدراسة (الفصول والمباحث)
# ─────────────────────────────────────────────────────────────────────
def make_plan(prs, data, T, chapters_data):
    slide = blank(prs)
    bg(slide, T["D"])
    oval(slide, W*0.60, -2, W*0.50, H*0.80, T["M"], 0.25)
    L   = get_layout(T)
    cy0 = L.section_dark(slide, T, "خطة الدراسة", "Plan de l'etude")
    n_ch = min(len(chapters_data), 4)
    if n_ch == 0:
        return slide
    gx   = 0.28
    cw   = min((W - 2.20 - gx*(n_ch-1)) / n_ch, 6.5)
    for ci, ch in enumerate(chapters_data[:4]):
        cx   = 1.10 + ci*(cw+gx)
        sc   = T["SC"][ci % len(T["SC"])]
        ch_h = H - cy0 - 0.35
        rect(slide, cx, cy0, cw, ch_h, T["M"])
        rect(slide, cx, cy0, cw, 0.06, sc)
        rect(slide, cx, cy0, 0.18, ch_h, sc)
        txt(slide, "F%d" % (ci+1), cx+0.24, cy0+0.08, cw-0.38, 0.58,
            font="Calibri", size=26, bold=True, color=sc, align=PP_ALIGN.RIGHT)
        secs = [s for s in ch.get("sections", []) if s][:5]
        title_h = 0.84 if secs else ch_h - 0.70
        txt(slide, safe(ch.get("title", "")), cx+0.24, cy0+0.70, cw-0.38, title_h,
            font=T["BF"], size=13, bold=True, color=T["TL"], align=PP_ALIGN.RIGHT)
        if secs:
            ln(slide, cx+0.24, cy0+1.58, cw-0.38, T["A"], 0.025)
            sec_h = (H - cy0 - 1.72 - 0.38) / len(secs)
            for si, sec in enumerate(secs):
                sy = cy0 + 1.72 + si*sec_h
                rect(slide, cx+0.30, sy+sec_h*0.42, 0.07, 0.07, sc)
                txt(slide, safe(sec), cx+0.44, sy+0.05, cw-0.58, sec_h-0.10,
                    font=T["BF"], size=11, color=T["TL"], align=PP_ALIGN.RIGHT)
    return slide


# ─────────────────────────────────────────────────────────────────────
# make_references — المراجع الرئيسية
# ─────────────────────────────────────────────────────────────────────
def make_references(prs, data, T):
    refs = [r for r in data.get("references", []) if r][:6]
    if not refs:
        return
    slide = blank(prs)
    bg(slide, T["L"])
    rect(slide, 0, 0, W, 0.18, T["A"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "ابرز المراجع", "Key References")
    avail = H - cy0 - 0.35
    rh    = min(1.10, avail / len(refs))
    for i, ref in enumerate(refs):
        ry  = cy0 + i*(rh+0.12)
        sc  = T["SC"][i % len(T["SC"])]
        bgc = T["CB"] if i % 2 == 0 else T["L"]
        rect(slide, 1.10, ry, W-2.20, rh, bgc)
        rect(slide, 1.10, ry, 0.18, rh, sc)
        txt(slide, str(i+1), 1.14, ry+0.06, 0.80, rh-0.12,
            font="Calibri", size=16, bold=True, color=sc, align=PP_ALIGN.CENTER)
        txt(slide, safe(ref), 2.00, ry+0.08, W-3.20, rh-0.16,
            font=T["BF"], size=11, color=T["TD"], align=PP_ALIGN.RIGHT)
    return slide


# ─────────────────────────────────────────────────────────────────────
# make_methodology_v2 — المنهجية الشاملة مع العينة والمجالات
# ─────────────────────────────────────────────────────────────────────
def make_methodology_v2(prs, data, T):
    slide = blank(prs)
    if T.get("LAYOUT") == "bold":
        rect(slide, 0, 0, 1.8, H, T["D"])
        rect(slide, 1.8, 0, W-1.8, H, T["L"])
    else:
        bg(slide, T["L"])
    rect(slide, 0, 0, W, 0.18, T["A"])
    L   = get_layout(T)
    cy0 = L.section_light(slide, T, "المنهجية والعينة", "Methodologie & Echantillon")

    meth     = safe(data.get("methodology", ""))
    dsource  = safe(data.get("dataSource", ""))
    stype    = safe(data.get("sampleType", ""))
    ssize    = safe(data.get("sampleSize", ""))
    tool_v   = safe(data.get("tool", ""))
    axes     = [a for a in data.get("toolAxes", []) if a]
    spatial  = safe(data.get("spatialScope", ""))
    temporal = safe(data.get("temporalScope", ""))
    human_s  = safe(data.get("humanScope", ""))
    software = safe(data.get("software", ""))
    tests    = [t_ for t_ in data.get("statisticalTests", []) if t_]

    boxes = []
    if meth:
        boxes.append(("المنهج المتبع", meth))
    if stype or ssize:
        boxes.append(("العينة", " | ".join(filter(None, [stype, ssize]))))
    if tool_v:
        tval = tool_v
        if axes:
            tval = tool_v + "\n" + "\n".join(["- " + a for a in axes[:4]])
        boxes.append(("ادوات الدراسة", tval))
    if spatial or temporal or human_s:
        scope_parts = []
        if spatial:  scope_parts.append("مكاني: " + spatial)
        if temporal: scope_parts.append("زماني: " + temporal)
        if human_s:  scope_parts.append("بشري: " + human_s)
        boxes.append(("مجالات الدراسة", "\n".join(scope_parts)))
    if software:
        sw = software
        if tests:
            sw = software + "\n" + " | ".join(tests[:4])
        boxes.append(("البرنامج والاختبارات", sw))
    if dsource:
        boxes.append(("مصدر البيانات", dsource))

    if not boxes:
        return slide

    n    = len(boxes)
    cols = min(n, 3)
    rows = math.ceil(n / cols)
    gx, gy = 0.28, 0.20
    avail_w = W - 2.20
    avail_h = H - cy0 - 0.35
    bw  = (avail_w - gx*(cols-1)) / cols
    bh  = (avail_h - gy*(rows-1)) / rows

    for i, (lbl, val) in enumerate(boxes):
        col  = i % cols
        row  = i // cols
        bx   = 1.10 + col*(bw+gx)
        by   = cy0 + row*(bh+gy)
        sc   = T["SC"][i % len(T["SC"])]
        light = T.get("LAYOUT") != "bold"
        bg_c = T["CB"] if light else T["M"]
        rect(slide, bx, by, bw, bh, bg_c)
        rect(slide, bx, by, bw, 0.055, sc)
        rect(slide, bx, by, 0.18, bh, sc)
        txt(slide, lbl, bx+0.26, by+0.08, bw-0.38, 0.52,
            font=T["BF"], size=13, bold=True,
            color=T["TD"] if light else T["TL"], align=PP_ALIGN.RIGHT)
        ln(slide, bx+0.26, by+0.64, bw-0.38, sc, 0.022)
        txt(slide, safe(val), bx+0.26, by+0.72, bw-0.38, max(bh-0.84, 0.40),
            font=T["BF"], size=11,
            color=T["TD"] if light else T["TL"], align=PP_ALIGN.RIGHT)
    return slide


# ═════════════════════════════════════════════════════════════════════
# MAIN ORCHESTRATOR
# ═════════════════════════════════════════════════════════════════════
def generate_presentation(data: dict, output_path: str) -> None:
    """
    البناء الرئيسي — يتحكم في الشرائح حسب slideStates + البيانات المدخلة
    هيكل الأستاذ حامدي ياسين: غلاف → مقدمة → خطة → إشكالية → منهجية → نتائج → خاتمة → شكر
    """
    key = data.get("theme", "navy_gold")
    T   = PALETTES.get(key, PALETTES["navy_gold"])
    prs = Presentation()
    prs.slide_width  = Cm(W)
    prs.slide_height = Cm(H)

    cfg   = data.get("slides", {})
    def show(k): return cfg.get(k, True)
    def fl(k):   return [x for x in data.get(k, []) if x]

    # 1. الغلاف — دائماً
    make_cover(prs, data, T)

    # 2. المقدمة
    if show("intro") and (data.get("introOverview") or data.get("introApproach")):
        make_intro(prs, data, T)

    # 3. خطة الدراسة
    chs = [c for c in data.get("chapters", []) if c.get("title")]
    if show("plan") and chs:
        make_plan(prs, data, T, chs)

    # 4. الإشكالية
    if show("problem") and (data.get("mainProblem") or data.get("mainQuestion") or fl("subQuestions")):
        make_problem(prs, data, T)

    # 5. الأهداف والفرضيات
    if show("objectives") and (fl("objectives") or fl("hypotheses")):
        make_objectives(prs, data, T)

    # 6. أهمية الدراسة
    if show("importance") and (fl("importance") or data.get("reasons")):
        make_importance_v2(prs, data, T)

    # 7. المنهجية والعينة
    if show("methodology") and (data.get("methodology") or data.get("sampleType") or data.get("tool")):
        make_methodology_v2(prs, data, T)

    # 8. لوحة KPI
    stats = [s for s in data.get("stats", []) if s.get("label") and s.get("value")]
    if show("kpi") and stats:
        make_stats(prs, data, T)

    # 9. النتائج
    if show("results") and fl("mainResults"):
        make_results(prs, data, T)

    # 10. الخاتمة
    if show("conclusion") and data.get("generalConclusion"):
        make_conclusion(prs, data, T)

    # 11. التوصيات
    if show("recommendations") and fl("recommendations"):
        make_recommendations(prs, data, T)

    # 12. الآفاق
    if show("future") and fl("futureWork"):
        make_future(prs, data, T)

    # 13. المراجع
    if show("references") and fl("references"):
        make_references(prs, data, T)

    # 14. شريحة الشكر
    if show("thankyou"):
        make_final(prs, data, T)

    prs.save(output_path)
    n = len(prs.slides._sldIdLst)
    print("✅  %d slides [%s·%s] → %s" % (n, T.get("LAYOUT","classic"), key, output_path),
          file=sys.stderr)

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python generator.py input.json output.pptx", file=sys.stderr)
        sys.exit(1)
    with open(sys.argv[1], encoding="utf-8") as f:
        payload = json.load(f)
    generate_presentation(payload, sys.argv[2])
